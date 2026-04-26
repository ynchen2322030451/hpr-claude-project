"""Create python-pptx shapes from parsed SVG elements.

Maps each SVG element to a native PPTX shape:
  SVGRect -> add_shape(RECTANGLE)
  SVGPath -> build_freeform() with line segments (cubics are linearized)
  SVGText -> add_textbox() with formatted run
"""

from __future__ import annotations

import logging
import math
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

from svg2pptx.path_parser import CmdType, PathSegment, linearize_cubic
from svg2pptx.style import ElementStyle, apply_fill, apply_line, hex_to_rgb
from svg2pptx.svg_parser import DrawableElement, SVGPath, SVGRect, SVGText
from svg2pptx.units import pt_to_emu

log = logging.getLogger(__name__)

# Minimum shape dimension in EMU (avoid degenerate shapes)
MIN_DIM = pt_to_emu(0.1)


class PPTXBuilder:
    """Build a PPTX presentation from parsed SVG elements."""

    def __init__(self, width_pt: float, height_pt: float):
        self.prs = Presentation()
        self.prs.slide_width = pt_to_emu(width_pt)
        self.prs.slide_height = pt_to_emu(height_pt)
        self._current_slide = None
        self._width_pt = width_pt
        self._height_pt = height_pt
        self._shape_count = 0
        self._skip_count = 0

    def new_slide(self) -> None:
        """Add a new blank slide."""
        layout = self.prs.slide_layouts[6]  # blank layout
        self._current_slide = self.prs.slides.add_slide(layout)

    @property
    def slide(self):
        if self._current_slide is None:
            self.new_slide()
        return self._current_slide

    def add_element(self, elem: DrawableElement) -> None:
        """Dispatch to the appropriate add method."""
        try:
            if isinstance(elem, SVGRect):
                self._add_rect(elem)
            elif isinstance(elem, SVGText):
                self._add_text(elem)
            elif isinstance(elem, SVGPath):
                self._add_path(elem)
            else:
                log.debug("Unknown element type: %s", type(elem))
        except Exception as e:
            log.debug("Failed to add element: %s", e)
            self._skip_count += 1

    def save(self, output_path: str) -> None:
        log.info("Saving PPTX: %d shapes, %d skipped -> %s",
                 self._shape_count, self._skip_count, output_path)
        self.prs.save(output_path)

    # ── Rect ────────────────────────────────────────────────────────────────

    def _add_rect(self, r: SVGRect) -> None:
        w_emu = max(pt_to_emu(r.width), MIN_DIM)
        h_emu = max(pt_to_emu(r.height), MIN_DIM)
        x_emu = pt_to_emu(r.x)
        y_emu = pt_to_emu(r.y)

        from pptx.enum.shapes import MSO_SHAPE
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_emu, y_emu, w_emu, h_emu
        )
        apply_fill(shape, r.style)
        apply_line(shape, r.style)
        self._shape_count += 1

    # ── Text ────────────────────────────────────────────────────────────────

    def _add_text(self, t: SVGText) -> None:
        info = t.info
        fs = max(info.font_size, 4.0)  # minimum readable

        # Estimate text box dimensions
        char_width_approx = fs * 0.6
        text_width = len(info.text) * char_width_approx + fs
        text_height = fs * 1.5

        x = info.x
        y = info.y - fs  # baseline to top adjustment

        is_rotated = abs(info.rotation) > 1.0

        if is_rotated:
            # For -90 degree (y-axis labels): swap width/height conceptually
            # and adjust position
            rot_deg = info.rotation
            if abs(rot_deg - 270.0) < 5.0 or abs(rot_deg + 90.0) < 5.0:
                rot_deg = -90.0
            # TextBox with rotation
            x_emu = pt_to_emu(x - text_height / 2)
            y_emu = pt_to_emu(y - text_width / 2)
            w_emu = max(pt_to_emu(text_width), MIN_DIM)
            h_emu = max(pt_to_emu(text_height), MIN_DIM)
        else:
            x_emu = pt_to_emu(x)
            y_emu = pt_to_emu(y)
            w_emu = max(pt_to_emu(text_width), MIN_DIM)
            h_emu = max(pt_to_emu(text_height), MIN_DIM)

        txbox = self.slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)

        if is_rotated:
            txbox.rotation = info.rotation if info.rotation > 0 else info.rotation + 360

        tf = txbox.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = info.text
        font = run.font
        font.size = Pt(fs)
        font.name = info.font_family
        font.bold = info.bold
        font.italic = info.italic

        if info.color and info.color != "none":
            try:
                font.color.rgb = hex_to_rgb(info.color)
            except Exception:
                pass

        self._shape_count += 1

    # ── Freeform path ───────────────────────────────────────────────────────

    def _is_circle_path(self, segments: List[PathSegment]
                        ) -> Optional[Tuple[float, float, float, float]]:
        """Detect M-C-C-C-C-Z circle pattern (4 cubic bezier quarter arcs).

        Returns (x, y, w, h) bounding box if detected, else None.
        """
        cmds = [s.cmd for s in segments]
        if cmds != [CmdType.MOVE, CmdType.CUBIC, CmdType.CUBIC,
                     CmdType.CUBIC, CmdType.CUBIC, CmdType.CLOSE]:
            return None

        # Collect all points to find bounding box
        all_pts = []
        for seg in segments:
            all_pts.extend(seg.points)
        if not all_pts:
            return None

        xs = [p[0] for p in all_pts]
        ys = [p[1] for p in all_pts]
        x0, x1 = min(xs), max(xs)
        y0, y1 = min(ys), max(ys)
        w, h = x1 - x0, y1 - y0

        # Check aspect ratio is roughly circular/elliptical (allow 30% deviation)
        if w < 0.01 or h < 0.01:
            return None
        ratio = w / h
        if ratio < 0.7 or ratio > 1.3:
            # Could still be an ellipse, use oval anyway
            pass

        return (x0, y0, w, h)

    def _add_path(self, p: SVGPath) -> None:
        segments = p.segments
        if not segments:
            return

        # Filter invisible paths
        style = p.style
        has_fill = style.fill and style.fill.lower() != "none"
        has_stroke = style.stroke and style.stroke.lower() != "none"
        if not has_fill and not has_stroke:
            return

        # Fast path: detect circles/ellipses and use OVAL shape
        circle_bbox = self._is_circle_path(segments)
        if circle_bbox is not None:
            x, y, w, h = circle_bbox
            from pptx.enum.shapes import MSO_SHAPE
            shape = self.slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                pt_to_emu(x), pt_to_emu(y),
                max(pt_to_emu(w), MIN_DIM), max(pt_to_emu(h), MIN_DIM),
            )
            apply_fill(shape, style)
            apply_line(shape, style)
            self._shape_count += 1
            return

        # Collect all points to find bounding box for the freeform builder
        all_pts = self._collect_line_points(segments)
        if len(all_pts) < 2:
            return

        # Build freeform shape
        start_x = pt_to_emu(all_pts[0][0])
        start_y = pt_to_emu(all_pts[0][1])

        try:
            fb = self.slide.shapes.build_freeform(start_x, start_y)
        except Exception:
            # Fallback: some pptx versions use different API
            log.debug("build_freeform failed, skipping path")
            self._skip_count += 1
            return

        for pt_coord in all_pts[1:]:
            fb.add_line_segments([(pt_to_emu(pt_coord[0]), pt_to_emu(pt_coord[1]))])

        try:
            shape = fb.convert_to_shape()
        except Exception as e:
            log.debug("convert_to_shape failed: %s", e)
            self._skip_count += 1
            return

        apply_fill(shape, style)
        apply_line(shape, style)
        self._shape_count += 1

    def _collect_line_points(self, segments: List[PathSegment]) -> List[Tuple[float, float]]:
        """Convert segments (with cubics linearized) into a flat point list."""
        pts: List[Tuple[float, float]] = []
        cx, cy = 0.0, 0.0
        subpath_start: Optional[Tuple[float, float]] = None

        for seg in segments:
            if seg.cmd == CmdType.MOVE:
                if seg.points:
                    cx, cy = seg.points[0]
                    subpath_start = (cx, cy)
                    pts.append((cx, cy))

            elif seg.cmd == CmdType.LINE:
                if seg.points:
                    cx, cy = seg.points[0]
                    pts.append((cx, cy))

            elif seg.cmd == CmdType.CUBIC:
                if len(seg.points) == 3:
                    p0 = (cx, cy)
                    p1, p2, p3 = seg.points
                    linearized = linearize_cubic(p0, p1, p2, p3, tolerance=0.3)
                    pts.extend(linearized)
                    cx, cy = seg.points[-1]

            elif seg.cmd == CmdType.CLOSE:
                # Close by going back to the subpath start (not first point overall)
                if subpath_start and pts:
                    pts.append(subpath_start)

        return pts
