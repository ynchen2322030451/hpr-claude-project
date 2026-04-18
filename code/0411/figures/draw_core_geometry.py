#!/usr/bin/env python3
"""
Draw MegaPower reactor core — tall-cylinder cutaway style matching Figure 2
in /参考几何结构图/Figure 2.pdf.

Key ideas
---------
* Tall cylinder (H ≈ 2.5×R); the figure looks like a reactor vessel.
* A 60° wedge is cut away from the main cylinder, revealing the two
  radial cut planes on the LEFT panel.  The RIGHT panel is the 60° wedge
  itself, viewed from a rotated azimuth so both of its faces are visible.
* Cut face shows axial zones (top→bottom):
      thin gold cap (upper monolith)
      large salmon/red fuel slab (≈ 75% of height)
      pin-pitch yellow vertical stripes = heat-pipe cross-sections
      thin green absorber band
      pink + off-white base bands
  Inner control-rod-channel column is pale pink + same bottom bands.
  Outer reflector column is gold, full height.
* Top face = full 2D top-down cross-section — drums, fuel/HP pins,
  channel, reflector (same geometry and colors as core_section_2d).
* Orthographic projection, painter's algorithm back-to-front.
* Cylinder wall is split into many vertical strips so depth sorts
  per-strip (not averaged over the whole wall).
* Cut planes are back-face culled by default.
"""

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import Polygon as MplPoly
from pathlib import Path

# ===================== Radial geometry =====================
WALL_1 = 8.84;   WALL_2 = 29.625;   R_REFL = 55.625
HP_D   = 1.575;  FUEL_DO = 1.425;   P_FUEL = 1.6
P_HP   = 1.6 * np.sqrt(3)
DRUM_CR = (WALL_2 + R_REFL) / 2
DRUM_R  = (R_REFL - WALL_2) / 2 - 0.4

# ===================== Axial zones =====================
# Axial heights mapped from OpenMC geometry (MEGA_OpenMC_test.py:165-170
# and generater.py:380-385):
#   H_core = 150 cm, H_reflector = 15 cm, H_gas_plenum = 20 cm
#   reflector_top  =  +90 cm        → Z_TOP      =    0.0
#   fuel_top       =  +75 cm        → Z_FUEL_TOP =  -15.0
#   fuel_bottom    =  -75 cm        → Z_FUEL_BOT = -165.0
#   reflector_bot  =  -90 cm        → Z_ABS_BOT  = -180.0   (axial reflector/BeO band)
#   reactor_bottom = -110 cm        → Z_BOT      = -200.0   (gas plenum base)
# Proportions therefore match the real reactor vessel 1:1.
Z_TOP       =    0.0
Z_FUEL_TOP  =  -15.0    # top BeO reflector above this line: HP + structure only
Z_FUEL_BOT  = -165.0
Z_ABS_TOP   = Z_FUEL_BOT
Z_ABS_BOT   = -180.0    # bottom BeO reflector (15 cm, analogous to top)
Z_PINK_BOT  = -200.0    # gas plenum (20 cm)
Z_BOT       = Z_PINK_BOT

# ===================== Colors =====================
# Major palette overhaul.  Design goals the earlier palette missed:
#   1. Hues must be clearly distinguishable: fuel-red / HP-yellow /
#      channel-pink used to blend together — they are now pulled apart
#      in both HUE and LUMINANCE (fuel → very dark maroon; HP → warm
#      straw gold; channel → cool lilac-pink; bottom band → dusty mauve).
#   2. Absorber green must not fight the gold/red family → switched from
#      vivid emerald to a muted deep PINE-TEAL that sits quietly in the
#      palette.
#   3. Reflector & drum must differ enough that drums pop → reflector is
#      brushed BRONZE (warm, medium-dark); drum is PALE-CHAMPAGNE (lighter
#      and more yellow) — the drums read clearly even from a distance.
#
# Coordinated "autumn museum" palette — all tones share a shared
# warm-earth family but with controlled hue spread.
C_REFL      = np.array([176, 140,  76]) / 255   # satin bronze
C_MONO      = np.array([155,  70,  62]) / 255   # oxidized brick (top face)
C_MONO_HI   = np.array([186,  95,  80]) / 255   # brick highlight (for gradient)
C_CHANNEL   = np.array([226, 210, 216]) / 255   # cool lilac-rose (channel top)
C_HP        = np.array([236, 198,  92]) / 255   # straw gold (HP circles)
C_FUEL      = np.array([ 86,  24,  30]) / 255   # very dark maroon fuel pellet
C_DRUM      = np.array([224, 192, 104]) / 255   # pale champagne drum
C_B4C       = np.array([ 44, 112,  94]) / 255   # deep pine-teal absorber

# Brushed metallic sheen — 3-stop bronze gradient, desaturated
C_SIDE_LIGHT = np.array([204, 170,  92]) / 255  # bronze highlight
C_SIDE_MID   = np.array([150, 116,  50]) / 255  # bronze base
C_SIDE_DARK  = np.array([ 92,  68,  26]) / 255  # bronze shadow

C_FUEL_SLAB  = np.array([232, 138, 128]) / 255  # monolith cut-face — SALMON
                                                # matching reference Fig. 2:
                                                # the monolith body is shown
                                                # as a flat salmon fill with no
                                                # fuel/HP internals visible.
C_TOPREFL    = np.array([232, 138, 128]) / 255  # top band background (same
                                                # salmon; HP pins drawn on
                                                # top as vertical lines).
                                                # (warm putty-tan — distinctly
                                                # different from the fuel slab
                                                # so you can clearly see that
                                                # above Z_FUEL_TOP the fuel has
                                                # ended but HP continues)
C_HP_STRIPE  = np.array([232, 196,  92]) / 255
C_CHAN_COL   = np.array([232, 218, 222]) / 255  # pale cool lilac channel
C_REFL_COL   = np.array([186, 148,  72]) / 255
C_ABS_BAND   = np.array([ 52, 122, 100]) / 255  # deep pine-teal band
C_PINK_BAND  = np.array([178, 132, 140]) / 255  # dusty mauve — clearly not
                                                # channel-pink and not fuel-red
C_RIM_DARK   = np.array([130,  96,  42]) / 255  # subtle bronze rim


# ===================== View (set per panel) =====================
_VIEW = {'azim': 215.0, 'elev': 25.0, 'zscale': 1.0,
         'right': None, 'up': None, 'fwd': None}


def set_view(azim, elev, zscale=1.0):
    a = np.deg2rad(azim); e = np.deg2rad(elev)
    fwd   = np.array([np.cos(a)*np.cos(e), np.sin(a)*np.cos(e), np.sin(e)])
    right = np.array([-np.sin(a), np.cos(a), 0.])
    up    = np.cross(fwd, right)
    _VIEW.update(azim=azim, elev=elev, zscale=zscale,
                 fwd=fwd, right=right, up=up)


def proj(pts):
    p = np.atleast_2d(pts).astype(float).copy()
    p[:, 2] *= _VIEW['zscale']
    return p @ _VIEW['right'], p @ _VIEW['up'], p @ _VIEW['fwd']


def shade(c, f):
    return np.clip(c * f, 0, 1)


def hex_dist(a, t):
    """Radial distance from origin to a regular hexagon boundary.

    Hex orientation: vertices at 30°, 90°, 150°, 210°, 270°, 330°; apothem
    directions (edge midpoints) at 0°, 60°, 120°, 180°, 240°, 300°.  Given
    apothem `a`, circumradius is 2a/√3.

    Old formula `a / cos(t % π/3 - π/6)` returns apothem at vertex angles
    (30°, 90°, 150°, 270°, …) — wrong for this orientation. The corrected
    shift by +π/6 puts the zeros of the cosine exactly on apothem rays so
    that f(0)=a, f(π/6)=2a/√3, f(π/2)=2a/√3, etc.
    """
    return a / np.cos((t + np.pi/6) % (np.pi/3) - np.pi/6)


# ===================== Pin positions =====================
def gen_hp():
    p = []
    for i in range(4, 13):
        X = np.arange(WALL_1 + HP_D/2, 1000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_HP/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def gen_fuel():
    p = []
    for i in range(3, 11):
        X = np.arange(WALL_1 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5 + P_FUEL*np.sqrt(3)/2, -1000,
                      -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-3], Y[i-3] + j*P_HP])
    for i in range(4, 12):
        X = np.arange(WALL_1 + P_FUEL/2 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def replicate(pts, n=6):
    o = []
    for k in range(n):
        a = np.deg2rad(60*k)
        c, s = np.cos(a), np.sin(a)
        o.append(pts @ np.array([[c, -s], [s, c]]).T)
    return np.vstack(o)


def drum_centers(n=6):
    return [np.array([DRUM_CR*np.cos(np.deg2rad(60*k)),
                      DRUM_CR*np.sin(np.deg2rad(60*k))]) for k in range(n)]


HP_ALL = replicate(gen_hp())
FUEL_ALL = replicate(gen_fuel())


# ===================== Render queue =====================
_FACE_QUEUE = []


def _push(depth_key, verts2d, fc, ec, lw):
    _FACE_QUEUE.append((depth_key, verts2d, fc, ec, lw))


def add_quad_3d(pts_xyz, fc, ec=None, lw=0.25, bias=0.0):
    v = np.array(pts_xyz)
    sx, sy, d = proj(v)
    verts2d = np.column_stack([sx, sy])
    if ec is None:
        ec = fc
    _push(d.mean() + bias, verts2d, fc, ec, lw)


def add_top_face_poly(xy_poly, z, fc, ec=None, lw=0.3, bias=0.0):
    v = np.array([[p[0], p[1], z] for p in xy_poly])
    sx, sy, _ = proj(v)
    verts2d = np.column_stack([sx, sy])
    if ec is None:
        ec = fc
    _push(1e6 + bias, verts2d, fc, ec, lw)


# ===================== Top face =====================
def in_visible(t, cut_a, cut_b):
    t = t % (2*np.pi)
    a = cut_a % (2*np.pi); b = cut_b % (2*np.pi)
    if a < b:
        return not (a <= t <= b)
    return not (t >= a or t <= b)


def top_circle(center, r, z, fc, ec=None, lw=0.3, bias=0, n=28):
    angs = np.linspace(0, 2*np.pi, n)
    poly = [[center[0] + r*np.cos(a), center[1] + r*np.sin(a)] for a in angs]
    add_top_face_poly(poly, z, fc, ec=ec if ec is not None else fc, lw=lw, bias=bias)


def _subtract_wedge(poly_xy, ca, cb):
    """Remove the pie-wedge (origin vertex, from angle ca CCW to angle cb,
    0<cb-ca<π) from the closed polygon poly_xy. Returns new polygon list
    of (x, y) tuples, closed (last == first)."""
    delta = (cb - ca) % (2*np.pi)
    if delta <= 1e-9 or delta >= np.pi - 1e-6:
        return [tuple(p) for p in poly_xy]

    def in_wedge(pt):
        if pt[0]*pt[0] + pt[1]*pt[1] < 1e-12:
            return False
        t = (np.arctan2(pt[1], pt[0]) - ca) % (2*np.pi)
        return 1e-7 < t < delta - 1e-7

    def intersect_ray(p0, p1, ang):
        c, s_ = np.cos(ang), np.sin(ang)
        denom = (p1[0]-p0[0])*s_ - (p1[1]-p0[1])*c
        if abs(denom) < 1e-12:
            return None
        s = (p0[1]*c - p0[0]*s_) / denom
        if s < 1e-7 or s > 1 - 1e-7:
            return None
        x = p0[0] + s*(p1[0]-p0[0])
        y = p0[1] + s*(p1[1]-p0[1])
        if x*c + y*s_ < 1e-7:
            return None
        return (s, (x, y))

    pts = [tuple(p) for p in poly_xy]
    if len(pts) >= 2 and (abs(pts[0][0]-pts[-1][0]) + abs(pts[0][1]-pts[-1][1]) < 1e-9):
        pts = pts[:-1]
    n = len(pts)
    if n == 0:
        return []

    start = None
    for i in range(n):
        if not in_wedge(pts[i]):
            start = i
            break
    if start is None:
        return []
    pts = pts[start:] + pts[:start]

    result = []
    state = 'out'
    for i in range(n):
        p0, p1 = pts[i], pts[(i+1) % n]
        ints = []
        for ang in (ca, cb):
            r = intersect_ray(p0, p1, ang)
            if r is not None:
                ints.append(r)
        ints.sort(key=lambda x: x[0])
        if state == 'out':
            result.append(p0)
        for s, pt in ints:
            if state == 'out':
                result.append(pt)
                state = 'in'
            else:
                result.append((0.0, 0.0))
                result.append(pt)
                state = 'out'
    if result and (abs(result[0][0]-result[-1][0]) + abs(result[0][1]-result[-1][1]) > 1e-9):
        result.append(result[0])
    return result


def _hex_sector(apothem, arc_start_rad, arc_end_rad):
    """Intersection of the regular hex (apothem=apothem, vertices at
    30°, 90°, …, 330°) with the pie-wedge from arc_start CCW to arc_end.
    Works for any arc size > 0 up to (but not including) 2π."""
    delta = (arc_end_rad - arc_start_rad) % (2*np.pi)
    if delta <= 1e-9:
        return []
    if delta >= 2*np.pi - 1e-6:
        hex_v = np.deg2rad([30, 90, 150, 210, 270, 330, 30])
        return [[hex_dist(apothem, t)*np.cos(t),
                 hex_dist(apothem, t)*np.sin(t)] for t in hex_v]
    R = apothem * 2.0 / np.sqrt(3)           # circum-radius
    out = [(0.0, 0.0)]
    r1 = hex_dist(apothem, arc_start_rad)
    out.append((r1*np.cos(arc_start_rad), r1*np.sin(arc_start_rad)))
    hex_angs = np.deg2rad([30, 90, 150, 210, 270, 330])
    rel = [(a - arc_start_rad) % (2*np.pi) for a in hex_angs]
    for i in sorted(range(6), key=lambda k: rel[k]):
        if 1e-7 < rel[i] < delta - 1e-7:
            a = hex_angs[i]
            out.append((R*np.cos(a), R*np.sin(a)))
    r2 = hex_dist(apothem, arc_end_rad)
    out.append((r2*np.cos(arc_end_rad), r2*np.sin(arc_end_rad)))
    out.append(out[0])
    return out


def draw_top_face(cut_a_deg, cut_b_deg, keep_wedge=False):
    """Top face at z = Z_TOP.
    cut_a_deg, cut_b_deg define the 60° wedge (ca CCW to cb).
    keep_wedge=False → render the 5/6 piece (everything *outside* the wedge).
    keep_wedge=True  → render the 1/6 piece (the wedge itself)."""
    ca = np.deg2rad(cut_a_deg); cb = np.deg2rad(cut_b_deg)
    wedge_delta = (cb - ca) % (2*np.pi)
    if keep_wedge:
        arc_start, arc_delta = ca, wedge_delta
    else:
        arc_start, arc_delta = cb, 2*np.pi - wedge_delta
    arc_end = arc_start + arc_delta

    margin = np.deg2rad(0.3)
    n_arc = max(24, int(arc_delta * 60 / np.pi))
    arc = np.linspace(arc_start + margin, arc_end - margin, n_arc)

    # Reflector disc sector
    refl_poly = [[0, 0]] + [[R_REFL*np.cos(a), R_REFL*np.sin(a)] for a in arc] + [[0, 0]]
    add_top_face_poly(refl_poly, Z_TOP, C_REFL, ec=shade(C_REFL, 0.78),
                      lw=0.5, bias=5)

    # Soft bronze rim — a very thin, only slightly darker ring inside the
    # reflector outer edge (NOT the earlier near-black ring, which looked
    # like a drawn outline).  Just enough to imply plate thickness.
    rim_R_out = R_REFL
    rim_R_in  = R_REFL * 0.992
    rim_poly_out = [[rim_R_out*np.cos(a), rim_R_out*np.sin(a)] for a in arc]
    rim_poly_in  = [[rim_R_in *np.cos(a), rim_R_in *np.sin(a)] for a in arc[::-1]]
    add_top_face_poly(rim_poly_out + rim_poly_in, Z_TOP,
                      shade(C_REFL, 0.82), ec=None,
                      lw=0.0, bias=5.3)

    # Hex monolith (outer), clipped to the sector — solid base layer.
    outer = _hex_sector(WALL_2, arc_start, arc_end)
    if outer:
        add_top_face_poly(outer, Z_TOP, C_MONO,
                          ec=shade(C_MONO, 0.60), lw=0.6, bias=6)
        # Soft radial shading: 2 concentric shrunk hex overlays, each
        # slightly brighter, to simulate a highlight toward the centre
        # (cheap alternative to a true raytraced gradient).
        for k, (scale, bright) in enumerate([(0.78, 1.08), (0.50, 1.15)]):
            shrunk = _hex_sector(WALL_2 * scale, arc_start, arc_end)
            if shrunk:
                add_top_face_poly(shrunk, Z_TOP, shade(C_MONO, bright),
                                  ec=None, lw=0.0, bias=6.05 + 0.02*k)

    # Inner hex channel, clipped to the sector
    inner = _hex_sector(WALL_1, arc_start, arc_end)
    if inner:
        add_top_face_poly(inner, Z_TOP, C_CHANNEL,
                          ec=shade(C_CHANNEL, 0.75), lw=0.5, bias=7)

    def in_sec(t):
        rel = (t - arc_start) % (2*np.pi)
        return 1e-7 < rel < arc_delta - 1e-7

    for p in FUEL_ALL:
        t = np.arctan2(p[1], p[0])
        if not in_sec(t):
            continue
        top_circle(p, FUEL_DO/2, Z_TOP, C_FUEL, ec='#6a1e16', lw=0.15, bias=8)
    for p in HP_ALL:
        t = np.arctan2(p[1], p[0])
        if not in_sec(t):
            continue
        top_circle(p, HP_D/2, Z_TOP, C_HP, ec='#8a6a10', lw=0.15, bias=8)

    for dc in drum_centers():
        tdc = np.arctan2(dc[1], dc[0])
        dang = np.arcsin(min(1, DRUM_R / np.linalg.norm(dc)))
        pad = dang + np.deg2rad(2)
        if not (in_sec(tdc) and in_sec(tdc - pad) and in_sec(tdc + pad)):
            continue
        top_circle(dc, DRUM_R, Z_TOP, shade(C_DRUM, 1.05),
                   ec=shade(C_DRUM, 0.65), lw=0.5, bias=9, n=72)
        # OpenMC geometry (MEGA_OpenMC_test.py:179-180, controldrum_angle=0):
        # drum_inner is offset INWARD from drum_outer by 2 cm, so the B4C
        # crescent (inside outer, outside inner) sits on the OUTWARD face
        # — away from the core. This matches the reference Figure 2.
        outward = np.arctan2(dc[1], dc[0])
        arc_a = np.linspace(outward - np.pi/3, outward + np.pi/3, 40)
        rdc_in = DRUM_R * 0.78
        crescent = [[dc[0]+DRUM_R*np.cos(a), dc[1]+DRUM_R*np.sin(a)] for a in arc_a]
        crescent += [[dc[0]+rdc_in*np.cos(a), dc[1]+rdc_in*np.sin(a)]
                     for a in arc_a[::-1]]
        add_top_face_poly(crescent, Z_TOP, C_B4C, ec='#1c6a30', lw=0.4, bias=10)


# ===================== Cylinder wall =====================
def draw_cylinder_side(cut_a_deg, cut_b_deg, n_strips=80, keep_wedge=False):
    ca = np.deg2rad(cut_a_deg); cb = np.deg2rad(cut_b_deg)
    wedge_delta = (cb - ca) % (2*np.pi)
    if keep_wedge:
        vs = ca
        ve = ca + wedge_delta
        # proportionally fewer strips for the smaller 60° arc (still smooth)
        n_strips = max(12, int(round(n_strips * wedge_delta / (2*np.pi))) + 6)
    else:
        vs = cb
        ve = cb + (2*np.pi - wedge_delta)
    margin = np.deg2rad(0.25)
    arc = np.linspace(vs + margin, ve - margin, n_strips + 1)

    view_az = np.deg2rad(_VIEW['azim'])
    for i in range(n_strips):
        a0, a1 = arc[i], arc[i+1]
        am = 0.5*(a0 + a1)
        lit = np.cos(am - view_az)          # -1..1 (back..front)
        t   = 0.5 * (1 + lit)               # 0..1

        # 3-stop metallic gradient: dark → mid → highlight → mid
        # with a narrow specular bump near the front facing strip.
        specular = np.exp(-((lit - 0.82) ** 2) / 0.012) * 0.35
        if t < 0.5:
            # blend dark → mid
            k = t / 0.5
            color = C_SIDE_DARK * (1 - k) + C_SIDE_MID * k
        else:
            # blend mid → light
            k = (t - 0.5) / 0.5
            color = C_SIDE_MID * (1 - k) + C_SIDE_LIGHT * k
        color = np.clip(color + specular, 0, 1)

        pts = [[R_REFL*np.cos(a0), R_REFL*np.sin(a0), Z_TOP],
               [R_REFL*np.cos(a1), R_REFL*np.sin(a1), Z_TOP],
               [R_REFL*np.cos(a1), R_REFL*np.sin(a1), Z_BOT],
               [R_REFL*np.cos(a0), R_REFL*np.sin(a0), Z_BOT]]
        v = np.array(pts)
        sx, sy, d = proj(v)
        verts2d = np.column_stack([sx, sy])
        # Edge same color, lw=0 → no visible vertical stripe lines
        _push(d.mean(), verts2d, color, color, 0.0)


# ===================== Bottom face =====================
def draw_bottom_face(cut_a_deg, cut_b_deg, keep_wedge=False):
    ca = np.deg2rad(cut_a_deg); cb = np.deg2rad(cut_b_deg)
    wedge_delta = (cb - ca) % (2*np.pi)
    if keep_wedge:
        vs = ca
        ve = ca + wedge_delta
    else:
        vs = cb
        ve = cb + (2*np.pi - wedge_delta)
    margin = np.deg2rad(0.3)
    arc_n = max(24, int((ve - vs) * 60 / np.pi))
    arc = np.linspace(vs + margin, ve - margin, arc_n)
    poly = [[0, 0, Z_BOT]] + [[R_REFL*np.cos(a), R_REFL*np.sin(a), Z_BOT] for a in arc]
    v = np.array(poly)
    sx, sy, _ = proj(v)
    verts2d = np.column_stack([sx, sy])
    _push(-1e5, verts2d, shade(C_SIDE_DARK, 0.95),
          shade(C_SIDE_DARK, 0.7), 0.3)


# ===================== Cut plane =====================
def cut_plane_visible(theta, wedge_center, is_wedge=False):
    """Back-face test for a radial cut plane.
    By default the plane belongs to the 5/6 piece (main cylinder with
    wedge removed); its exposed side points INTO the wedge (toward
    wedge_center). For the isolated wedge piece, the exposed side
    points OUTWARD — away from wedge_center — so invert."""
    d = (wedge_center - theta + np.pi) % (2*np.pi) - np.pi
    normal_ang = theta + (np.pi/2 if d > 0 else -np.pi/2)
    if is_wedge:
        normal_ang += np.pi
    view = np.deg2rad(_VIEW['azim'])
    return np.cos(normal_ang - view) > 0


def draw_cut_plane(theta, wedge_center, *, force=False, bias_offset=0.0,
                   is_wedge=False):
    """Cut plane at angle theta, exposed side toward wedge_center.
    If back-face culling would skip it and force is False, skip it.
    `bias_offset` lets the caller nudge one plane slightly forward so
    that, for the isolated wedge panel, the viewer-closer plane wins."""
    if not force and not cut_plane_visible(theta, wedge_center, is_wedge=is_wedge):
        return

    c, s = np.cos(theta), np.sin(theta)

    def quad(r1, r2, z1, z2, color, bias=0, ec=None, lw=0.2):
        pts = [[r1*c, r1*s, z1],
               [r2*c, r2*s, z1],
               [r2*c, r2*s, z2],
               [r1*c, r1*s, z2]]
        add_quad_3d(pts, color,
                    ec=ec if ec is not None else shade(np.asarray(color), 0.72),
                    lw=lw, bias=bias + bias_offset)

    # --- Metallic gradient: shade differs for the "lit" vs "shadowed"
    # cut plane so the two cut faces don't look flat-identical.
    view = np.deg2rad(_VIEW['azim'])
    d = (wedge_center - theta + np.pi) % (2*np.pi) - np.pi
    normal_ang = theta + (np.pi/2 if d > 0 else -np.pi/2)
    if is_wedge:
        normal_ang += np.pi
    face_lit = 0.5 * (1 + np.cos(normal_ang - view))          # 0..1
    gold_f   = 0.82 + 0.22 * face_lit                         # brass highlight
    slab_f   = 0.90 + 0.15 * face_lit                         # copper highlight

    # Hex radial boundaries at this cut angle.  The cut plane lies along
    # a hex vertex ray (θ ∈ {150°, 270°, 210°}), so `hex_dist` returns the
    # circumradius there — this is what makes the top-face hex polygon and
    # the vertical cut-plane wall share the SAME outer edge (fixing the
    # prior alignment mismatch where the cut wall stopped at apothem).
    r_outer = hex_dist(WALL_2, theta)
    r_inner = hex_dist(WALL_1, theta)

    # Reflector column (hex_outer → R_REFL, brushed bronze)
    quad(r_outer, R_REFL, Z_BOT, Z_TOP, shade(C_REFL_COL, gold_f), bias=-5)

    # Central channel column (0 → hex_inner).  Pale cool-lilac pink from
    # the TOP of the bottom mauve band up to the crown.
    quad(0, r_inner, Z_ABS_BOT, Z_TOP, C_CHAN_COL, bias=-5)

    # Monolith annulus (hex_inner → hex_outer).  Axial stack top→bottom,
    # matching the OpenMC geometry precisely (H_core=150, H_reflector=15):
    #   1. Z_FUEL_TOP → Z_TOP        top BeO reflector band (15 cm).
    #      Fuel has ended here but HPs continue up; background colour is
    #      the putty-tan "cap" tone so the user can see the "一根根" HP
    #      stripes protruding through a visibly different zone.
    #   2. Z_FUEL_BOT → Z_FUEL_TOP   active fuel slab (150 cm, brick).
    #   3. Z_ABS_BOT  → Z_FUEL_BOT   bottom BeO reflector band (teal-green
    #      cosmetic stripe in the reference figure; kept here to match).
    # Monolith cut face:
    #   * Body (Z_FUEL_BOT → Z_FUEL_TOP):  SOLID salmon, no internals.
    #   * Top band (Z_FUEL_TOP → Z_TOP):   NO slab — only HP pin stripes
    #     are drawn below, so the figure background shows through between
    #     rods (requested look: monolith does NOT fill around the rods
    #     above Z_FUEL_TOP).
    quad(r_inner, r_outer, Z_FUEL_BOT, Z_FUEL_TOP, shade(C_FUEL_SLAB, slab_f),  bias=-5)
    quad(r_inner, r_outer, Z_ABS_BOT,  Z_ABS_TOP,  C_ABS_BAND,                  bias=-5)

    # Bottom mauve band: spans the ENTIRE hex interior (0 → hex_outer),
    # representing the gas-plenum region below the bottom BeO reflector.
    quad(0, r_outer, Z_PINK_BOT, Z_ABS_BOT, C_PINK_BAND, bias=-5)

    # 3D-shaded pin bundles in the top band (Z_FUEL_TOP → Z_TOP).
    # Each rod is rendered as N_SUB vertical sub-stripes with a cosine
    # brightness profile across its diameter, giving a cylindrical
    # (3D) appearance.  Both HP rod bundles (gold, thicker) and fuel
    # rod bundles (maroon, thinner) are shown, interleaved, standing
    # alone against the figure background.
    MARGIN  = 2.0
    hp_hw   = HP_D / 2.0
    fuel_hw = FUEL_DO / 2.0
    r_hp = np.linspace(r_inner + MARGIN + hp_hw,
                       r_outer - MARGIN - hp_hw, 9)
    r_fuel = 0.5 * (r_hp[:-1] + r_hp[1:])
    N_SUB = 9  # sub-stripes per rod for cylindrical shading

    def draw_rod_3d(r_center, hw, base_color, bias):
        """Render a single rod as N_SUB vertical slices with cosine
        shading across its diameter so it reads as a 3D cylinder."""
        edges = np.linspace(-hw, +hw, N_SUB + 1)
        for i in range(N_SUB):
            u0, u1 = edges[i], edges[i+1]
            uc = 0.5 * (u0 + u1)
            # Cosine brightness: bright at centre, dark at rim
            bright = 0.50 + 0.65 * np.cos(0.5 * np.pi * uc / hw)
            col = shade(base_color, bright)
            # Edge-slice gets a darker stroke for silhouette
            is_rim = (i == 0) or (i == N_SUB - 1)
            ec = shade(base_color, 0.40) if is_rim else col
            lw = 0.25 if is_rim else 0.0
            quad(r_center + u0, r_center + u1,
                 Z_FUEL_TOP, Z_TOP,
                 col, bias=bias + i * 0.01,
                 ec=ec, lw=lw)

    # HP bundles — bright gold, full top-band height
    for r in r_hp:
        draw_rod_3d(r, hp_hw, C_HP_STRIPE, bias=500)
    # Fuel bundles — deep maroon, interleaved between HPs
    for r in r_fuel:
        draw_rod_3d(r, fuel_hw, C_FUEL, bias=495)


# ===================== Panel renderer =====================
def _wedge_center_rad(ca, cb):
    """Centre of the SHORT arc from ca to cb (in radians)."""
    dab = (cb - ca) % (2*np.pi)
    if dab <= np.pi:
        return ca + 0.5*dab
    return ca - 0.5*((ca - cb) % (2*np.pi))


def render_panel(ax, cut_a_deg, cut_b_deg, *, azim=215.0, elev=25.0,
                 is_wedge=False):
    """Render one cutaway panel.
    is_wedge=False → main cylinder with the wedge REMOVED (5/6 piece).
    is_wedge=True  → the 60° wedge PIECE alone."""
    global _FACE_QUEUE
    _FACE_QUEUE = []
    set_view(azim, elev)

    ca = np.deg2rad(cut_a_deg); cb = np.deg2rad(cut_b_deg)
    wedge_center = _wedge_center_rad(ca, cb)

    draw_bottom_face(cut_a_deg, cut_b_deg, keep_wedge=is_wedge)
    draw_cylinder_side(cut_a_deg, cut_b_deg, keep_wedge=is_wedge)
    draw_cut_plane(ca, wedge_center, is_wedge=is_wedge)
    draw_cut_plane(cb, wedge_center, is_wedge=is_wedge)
    draw_top_face(cut_a_deg, cut_b_deg, keep_wedge=is_wedge)

    _FACE_QUEUE.sort(key=lambda f: f[0])
    for _, v, fc, ec, lw in _FACE_QUEUE:
        ax.add_patch(MplPoly(v, closed=True, fc=fc, ec=ec, lw=lw))
    ax.set_aspect('equal')
    ax.axis('off')
    ax.autoscale_view()
    # Return a copy of the sorted face queue so the PPTX writer can emit
    # each colour patch as a separate draggable freeform shape.
    return list(_FACE_QUEUE)


# ===================== Figure =====================
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(20, 12),
                               gridspec_kw={'width_ratios': [1.6, 1.0],
                                            'wspace': 0.02})

ELEV = 25

# LEFT: 4/6 cylinder (keeping 240°). The removed 120° wedge spans two
# 1/6 bays centred on drums at 180° and 240°, so the cut planes at
# 150° and 270° fall on natural 1/6 boundaries (midway between drums).
LEFT_AZIM  = 210
LEFT_CA    = 150                  # cut plane A
LEFT_CB    = 270                  # cut plane B  (removed wedge = 120°)
left_queue = render_panel(ax1, LEFT_CA, LEFT_CB, azim=LEFT_AZIM, elev=ELEV,
                          is_wedge=False)

# RIGHT: an axisymmetric 1/6 wedge centred on a control drum (drum at
# 240° sits in the middle of the 210°–270° wedge). The view azimuth is
# mirrored across the current symmetry axis so the dominant cut face
# is on the opposite side compared to the previous build — i.e. the
# two cut faces "swap" their on-screen positions.
RIGHT_CA   = 210                  # 1/6 wedge CA
RIGHT_CB   = 270                  # 1/6 wedge CB  (wedge = 60°, drum at 240°)
# Mirror the panel so the wedge "points" to the LEFT (symmetry axis on the
# left-hand side of the figure) instead of to the right.  Azimuth =
# CA − 60° places the viewer on the opposite side of the wedge bisector.
RIGHT_AZIM = (RIGHT_CA - 60) % 360   # = 150°
right_queue = render_panel(ax2, RIGHT_CA, RIGHT_CB,
                           azim=RIGHT_AZIM, elev=ELEV, is_wedge=True)

plt.tight_layout(pad=1.0)

out = Path(__file__).parent
for fmt in ['svg', 'pdf', 'png', 'jpg']:
    p = out / f'core_cross_section_3d.{fmt}'
    fig.savefig(p, format='jpeg' if fmt == 'jpg' else fmt,
                dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    print(f"Saved: {p}")

# -------- PPTX wrapper (draggable freeform shapes, no bitmap) --------
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: F401


def _rgb(c):
    if isinstance(c, str):
        s = c.lstrip('#')
        if len(s) == 3:
            s = ''.join(ch*2 for ch in s)
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    r, g, b = [int(round(max(0, min(1, float(v))) * 255)) for v in c[:3]]
    return RGBColor(r, g, b)


def _panel_bbox(queue):
    all_v = np.concatenate([v for _, v, *_ in queue], axis=0)
    x0, y0 = all_v[:, 0].min(), all_v[:, 1].min()
    x1, y1 = all_v[:, 0].max(), all_v[:, 1].max()
    return x0, y0, x1, y1


def _draw_queue_as_freeforms(slide, queue, box_left_emu, box_top_emu,
                             box_w_emu, box_h_emu):
    """Emit every polygon in the (already depth-sorted) queue as a
    separate freeform shape. All shapes can be selected, re-coloured,
    dragged and regrouped in PowerPoint."""
    x0, y0, x1, y1 = _panel_bbox(queue)
    pw, ph = (x1 - x0), (y1 - y0)
    scale = min(box_w_emu / pw, box_h_emu / ph) * 0.97
    # Centre inside the box
    off_x = box_left_emu + (box_w_emu - pw * scale) / 2
    off_y = box_top_emu  + (box_h_emu - ph * scale) / 2

    def to_emu(v):
        # matplotlib y-axis is up; PPT y-axis is down → flip.
        ex = off_x + (v[0] - x0) * scale
        ey = off_y + (y1 - v[1]) * scale
        return Emu(int(ex)), Emu(int(ey))

    for _, verts, fc, ec, lw in queue:
        pts = [to_emu(v) for v in verts]
        if len(pts) < 3:
            continue
        start = pts[0]
        ff = slide.shapes.build_freeform(start[0], start[1], scale=1.0)
        segs = pts[1:]
        ff.add_line_segments(segs, close=True)
        shp = ff.convert_to_shape()
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fc)
        line = shp.line
        # Stroke off by default (each polygon already carries its own
        # fill colour). Flip to a visible border if the source lw>0.
        if lw and lw > 0:
            line.color.rgb = _rgb(ec)
            line.width = Pt(max(0.25, float(lw)))
        else:
            line.fill.background()


prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# White background rectangle (so the slide looks clean beneath the shapes)
margin = Inches(0.25)
total_w = prs.slide_width - 2 * margin
total_h = prs.slide_height - 2 * margin
# Match the figure's 1.6:1 column ratio so left panel is bigger
left_w  = int(total_w * (1.6 / 2.6))
right_w = int(total_w * (1.0 / 2.6))
gap     = Inches(0.15)
left_w_eff  = left_w - gap // 2
right_w_eff = right_w - gap // 2

_draw_queue_as_freeforms(slide, left_queue,
                         int(margin), int(margin),
                         int(left_w_eff), int(total_h))
_draw_queue_as_freeforms(slide, right_queue,
                         int(margin) + int(left_w) + int(gap // 2),
                         int(margin),
                         int(right_w_eff), int(total_h))

pptx_path = out / 'core_cross_section_3d.pptx'
prs.save(str(pptx_path))
print(f"Saved: {pptx_path}")
plt.close(fig)
