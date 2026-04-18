#!/usr/bin/env python3
"""
Pure 2D top-down cross-section of the MegaPower core.

Outputs:
  core_section_2d.svg / .pdf / .png
  core_section_2d.pptx        — single slide with two panels (full disk + 1/6 wedge)

This figure is intended to be used as a "flat" base drawing that can be manually
extruded in PowerPoint (Format Shape → 3-D Rotation + 3-D Format → Depth) to
obtain a stereoscopic appearance without any 3-D rendering code.
"""

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import Polygon as MplPoly, Wedge as MplWedge
from pathlib import Path

# --- Geometry (same as 3D script) ---
WALL_1 = 8.84;   WALL_2 = 29.625;   R_REFL = 55.625
HP_D   = 1.575;  FUEL_DO = 1.425;   P_FUEL = 1.6
P_HP   = 1.6 * np.sqrt(3)
DRUM_CR = (WALL_2 + R_REFL) / 2
DRUM_R  = (R_REFL - WALL_2) / 2 - 0.4

C_REFL    = np.array([218, 165,  32]) / 255
C_MONO    = np.array([218, 165,  32]) / 255
C_CHANNEL = np.array([250, 210, 210]) / 255
C_HP      = np.array([255, 215,  45]) / 255
C_FUEL    = np.array([192,  57,  43]) / 255
C_DRUM    = np.array([235, 190,  55]) / 255
C_B4C     = np.array([ 50, 205,  50]) / 255


def gen_hp():
    p = []
    for i in range(4, 13):
        X = np.arange(WALL_1 + HP_D/2, 1000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_HP/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def gen_fuel():
    p = []
    for i in range(3, 11):
        X = np.arange(WALL_1 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5 + P_FUEL*np.sqrt(3)/2, -1000, -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-3], Y[i-3] + j*P_HP])
    for i in range(4, 12):
        X = np.arange(WALL_1 + P_FUEL/2 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def replicate(pts, n=6):
    o = []
    for k in range(n):
        a = np.deg2rad(60*k)
        R = np.array([[np.cos(a), -np.sin(a)], [np.sin(a), np.cos(a)]])
        o.append(pts @ R.T)
    return np.vstack(o)


def drum_centers(n=6):
    return [np.array([DRUM_CR*np.cos(np.deg2rad(60*k)),
                      DRUM_CR*np.sin(np.deg2rad(60*k))]) for k in range(n)]


HP = replicate(gen_hp())
FUEL = replicate(gen_fuel())


def draw_hex(ax, r, **kw):
    hex_a = np.deg2rad([30, 90, 150, 210, 270, 330])
    verts = [[r/np.cos(np.pi/6)*np.cos(t), r/np.cos(np.pi/6)*np.sin(t)] for t in hex_a]
    ax.add_patch(MplPoly(verts, closed=True, **kw))


def draw_wedge_pie(ax, theta1_deg, theta2_deg):
    """Draw the full core cross-section as a wedge (theta1 to theta2)."""
    # Reflector
    ax.add_patch(MplWedge((0, 0), R_REFL, theta1_deg, theta2_deg,
                          fc=C_REFL, ec='#8a6715', lw=1.0, zorder=1))
    # Monolith hex (full hex, clipped visually by wedge via zorder stacking)
    if abs((theta2_deg - theta1_deg) % 360) in (0, 360):
        draw_hex(ax, WALL_2, fc=C_MONO, ec='#8a6715', lw=0.8, zorder=2)
    else:
        # Intersect hex with wedge sector: simple approach — draw full hex then
        # overlay the reflector wedge OUTSIDE the visible sector with same gold
        draw_hex(ax, WALL_2, fc=C_MONO, ec='#8a6715', lw=0.8, zorder=2)

    # Inner hex (channel)
    draw_hex(ax, WALL_1, fc=C_CHANNEL, ec='#8a6040', lw=0.8, zorder=10)

    # Pins
    def ang_ok(p):
        t = np.rad2deg(np.arctan2(p[1], p[0])) % 360
        a1, a2 = theta1_deg % 360, theta2_deg % 360
        if a1 < a2:
            return a1 <= t <= a2
        return t >= a1 or t <= a2

    for p in FUEL:
        if ang_ok(p):
            ax.add_patch(plt.Circle(p, FUEL_DO/2, fc=C_FUEL, ec='#6a1e16',
                                    lw=0.3, zorder=5))
    for p in HP:
        if ang_ok(p):
            ax.add_patch(plt.Circle(p, HP_D/2, fc=C_HP, ec='#8a6a10',
                                    lw=0.3, zorder=6))

    # Drums
    for dc in drum_centers():
        tdc_deg = np.rad2deg(np.arctan2(dc[1], dc[0])) % 360
        dang_deg = np.rad2deg(np.arcsin(min(1, DRUM_R/np.linalg.norm(dc))))
        # Require the whole drum to be inside the visible wedge
        def inside(t):
            a1, a2 = theta1_deg % 360, theta2_deg % 360
            tn = t % 360
            if a1 < a2:
                return a1 <= tn <= a2
            return tn >= a1 or tn <= a2
        if not (inside(tdc_deg) and inside(tdc_deg-dang_deg-2)
                and inside(tdc_deg+dang_deg+2)):
            continue
        ax.add_patch(plt.Circle(dc, DRUM_R, fc=C_DRUM, ec='#8a6715',
                                lw=0.8, zorder=7))
        # B4C crescent: annular band on inner side, 120°
        inward = np.arctan2(-dc[1], -dc[0])
        arc = np.linspace(inward - np.pi/3, inward + np.pi/3, 40)
        r_in = DRUM_R * 0.78
        pts = [[dc[0]+DRUM_R*np.cos(a), dc[1]+DRUM_R*np.sin(a)] for a in arc]
        pts += [[dc[0]+r_in*np.cos(a), dc[1]+r_in*np.sin(a)] for a in arc[::-1]]
        ax.add_patch(MplPoly(pts, closed=True, fc=C_B4C, ec='#1f6b1f',
                             lw=0.6, zorder=8))


# ================= Figure =================
fig, (axL, axR) = plt.subplots(1, 2, figsize=(16, 8),
                               gridspec_kw={'width_ratios': [1.4, 1.0]})

# Left: full disk (0 to 360)
draw_wedge_pie(axL, 0, 360)
axL.set_xlim(-R_REFL*1.08, R_REFL*1.08)
axL.set_ylim(-R_REFL*1.08, R_REFL*1.08)
axL.set_aspect('equal'); axL.axis('off')
axL.set_title('Full core cross-section', fontsize=14, pad=12)

# Right: 1/6 wedge (270° to 330°) — clip everything to wedge shape
draw_wedge_pie(axR, 270, 330)
# Clip all patches on axR to the 60° sector
clip_wedge = MplWedge((0, 0), R_REFL*1.05, 270, 330,
                      transform=axR.transData)
for patch in list(axR.patches):
    patch.set_clip_path(clip_wedge)
axR.set_xlim(-5, R_REFL*1.08)
axR.set_ylim(-R_REFL*1.08, R_REFL*0.2)
axR.set_aspect('equal'); axR.axis('off')
axR.set_title('1/6 symmetry sector', fontsize=14, pad=12)

plt.tight_layout(pad=1.5)

out = Path(__file__).parent
for fmt in ['svg', 'pdf', 'png']:
    p = out / f'core_section_2d.{fmt}'
    fig.savefig(p, format=fmt, dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    print(f"Saved: {p}")

# -------- PPTX (embed PNG, two-panel) --------
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])
iw, ih = Inches(12.5), Inches(6.8)
sl.shapes.add_picture(str(out / 'core_section_2d.png'),
                      (prs.slide_width - iw)//2,
                      (prs.slide_height - ih)//2, iw, ih)
pptx_path = out / 'core_section_2d.pptx'
prs.save(str(pptx_path))
print(f"Saved: {pptx_path}")
plt.close(fig)
