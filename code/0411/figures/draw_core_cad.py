#!/usr/bin/env python3
"""
CAD-render aesthetic of the MegaPower core.

This is the "CAD version" companion to draw_core_geometry.py:
the same geometry but rendered in an engineering / CAD-render style —
muted palette, thin black edge lines on every face, flat shading —
so it reads like a technical drawing rather than a stylised
illustration.

Outputs:  core_cad_render.{svg,pdf,png,jpg,pptx}
"""

import numpy as np
import matplotlib.pyplot as plt
from matplotlib.patches import Polygon as MplPoly
from pathlib import Path

# -------- Geometry constants (mirror draw_core_geometry.py) --------
WALL_1 = 8.84;   WALL_2 = 29.625;   R_REFL = 55.625
HP_D   = 1.575;  FUEL_DO = 1.425;   P_FUEL = 1.6
P_HP   = 1.6 * np.sqrt(3)
DRUM_CR = (WALL_2 + R_REFL) / 2
DRUM_R  = (R_REFL - WALL_2) / 2 - 0.4

# Z coordinates 1:1 mapped from OpenMC geometry (MEGA_OpenMC_test.py):
#   reflector_top  = +90 cm  → Z_TOP      =    0.0
#   fuel_top       = +75 cm  → Z_FUEL_TOP =  -15.0  (15 cm top BeO reflector)
#   fuel_bottom    = -75 cm  → Z_FUEL_BOT = -165.0  (150 cm active fuel)
#   (bottom BeO 15 cm)        → Z_ABS_BOT  = -180.0
#   reactor_bottom = -110 cm → Z_PINK_BOT = -200.0  (20 cm gas plenum)
Z_TOP      =    0.0
Z_FUEL_TOP =  -15.0
Z_FUEL_BOT = -165.0
Z_ABS_TOP  = Z_FUEL_BOT
Z_ABS_BOT  = -180.0
Z_PINK_BOT = -200.0
Z_BOT      = Z_PINK_BOT

# -------- CAD-render palette (muted / desaturated) -----------------
# Synced with draw_core_geometry.py (earth-tone / SCI journal style).
C_SIDE       = np.array([176, 140,  76]) / 255      # vessel wall — satin bronze
C_REFL       = np.array([176, 140,  76]) / 255
C_MONO       = np.array([176, 140,  76]) / 255      # top face reflector/monolith
C_CHANNEL    = np.array([226, 210, 216]) / 255      # cool lilac channel
C_TOPREFL    = np.array([232, 138, 128]) / 255      # top HP-stripe band bg
C_FUEL_SLAB  = np.array([232, 138, 128]) / 255      # monolith cut face — salmon
                                                     # (reference Fig. 2 style:
                                                     # body is flat salmon, no
                                                     # internals shown)
C_FUEL       = np.array([186,  82,  72]) / 255      # fuel rod tops on TOP face
C_HP_STRIPE  = np.array([232, 196,  92]) / 255      # HP rod — straw gold
C_ABS_BAND   = np.array([ 52, 122, 100]) / 255      # bottom BeO band — pine teal
C_PINK_BAND  = np.array([178, 132, 140]) / 255      # gas plenum — dusty mauve
C_DRUM       = np.array([224, 192, 104]) / 255      # drum face — pale champagne
C_B4C        = np.array([ 44, 112,  94]) / 255      # B4C crescent — pine teal

EDGE = '#1a1a1a'          # thin black edge on every face
LW_THIN  = 0.35
LW_MED   = 0.55
LW_HEAVY = 0.9

# -------- View ------------------------------------------------------
_VIEW = {'azim': 215.0, 'elev': 25.0, 'zscale': 1.0,
         'right': None, 'up': None, 'fwd': None}


def set_view(azim, elev, zscale=1.0):
    a = np.deg2rad(azim); e = np.deg2rad(elev)
    fwd   = np.array([np.cos(a)*np.cos(e), np.sin(a)*np.cos(e), np.sin(e)])
    right = np.array([-np.sin(a), np.cos(a), 0.])
    up    = np.cross(fwd, right)
    _VIEW.update(azim=azim, elev=elev, zscale=zscale,
                 fwd=fwd, right=right, up=up)


def proj(pts):
    p = np.atleast_2d(pts).astype(float).copy()
    p[:, 2] *= _VIEW['zscale']
    return p @ _VIEW['right'], p @ _VIEW['up'], p @ _VIEW['fwd']


def hex_dist(a, t):
    # Hex with vertices at 30°/90°/…/330° (matches hex_sector below).
    # The π/6 phase shift is critical: at cut angles 150°/270° (which
    # ARE hex vertex angles) this returns the circumradius, so the
    # top-face hex polygon and the vertical cut-plane wall share the
    # same outer edge — no notch.
    return a / np.cos((t + np.pi/6) % (np.pi / 3) - np.pi / 6)


# -------- Render queue ---------------------------------------------
_Q = []


def _push(d, v, fc, ec, lw):
    _Q.append((d, v, fc, ec, lw))


def quad3d(pts_xyz, fc, ec=EDGE, lw=LW_THIN, bias=0.0):
    v = np.array(pts_xyz)
    sx, sy, d = proj(v)
    _push(d.mean() + bias, np.column_stack([sx, sy]), fc, ec, lw)


def top_poly(xy, z, fc, ec=EDGE, lw=LW_MED, bias=0.0):
    v = np.array([[p[0], p[1], z] for p in xy])
    sx, sy, _ = proj(v)
    _push(1e6 + bias, np.column_stack([sx, sy]), fc, ec, lw)


# -------- Hex ∩ pie wedge ------------------------------------------
def hex_sector(apothem, arc_start, arc_end):
    delta = (arc_end - arc_start) % (2*np.pi)
    if delta <= 1e-9:
        return []
    if delta >= 2*np.pi - 1e-6:
        v = np.deg2rad([30, 90, 150, 210, 270, 330, 30])
        return [[hex_dist(apothem, t)*np.cos(t),
                 hex_dist(apothem, t)*np.sin(t)] for t in v]
    R = apothem * 2.0 / np.sqrt(3)
    out = [(0, 0)]
    r1 = hex_dist(apothem, arc_start)
    out.append((r1*np.cos(arc_start), r1*np.sin(arc_start)))
    hex_angs = np.deg2rad([30, 90, 150, 210, 270, 330])
    rel = [(a - arc_start) % (2*np.pi) for a in hex_angs]
    for i in sorted(range(6), key=lambda k: rel[k]):
        if 1e-7 < rel[i] < delta - 1e-7:
            a = hex_angs[i]
            out.append((R*np.cos(a), R*np.sin(a)))
    r2 = hex_dist(apothem, arc_end)
    out.append((r2*np.cos(arc_end), r2*np.sin(arc_end)))
    out.append(out[0])
    return out


def circle_poly(cx, cy, r, n=48):
    t = np.linspace(0, 2*np.pi, n)
    return [[cx + r*np.cos(a), cy + r*np.sin(a)] for a in t]


# -------- Top face --------------------------------------------------
def draw_top_face(ca_deg, cb_deg, keep_wedge=False):
    ca = np.deg2rad(ca_deg); cb = np.deg2rad(cb_deg)
    wd = (cb - ca) % (2*np.pi)
    if keep_wedge:
        a0, ad = ca, wd
    else:
        a0, ad = cb, 2*np.pi - wd
    a1 = a0 + ad
    margin = np.deg2rad(0.3)
    n = max(24, int(ad * 60 / np.pi))
    arc = np.linspace(a0 + margin, a1 - margin, n)

    # Reflector sector
    refl = [[0, 0]] + [[R_REFL*np.cos(a), R_REFL*np.sin(a)] for a in arc] + [[0, 0]]
    top_poly(refl, Z_TOP, C_REFL, ec=EDGE, lw=LW_HEAVY, bias=5)

    # Hex monolith + inner channel, clipped to sector
    outer = hex_sector(WALL_2, a0, a1)
    if outer:
        top_poly(outer, Z_TOP, C_MONO, ec=EDGE, lw=LW_MED, bias=6)
    inner = hex_sector(WALL_1, a0, a1)
    if inner:
        top_poly(inner, Z_TOP, C_CHANNEL, ec=EDGE, lw=LW_MED, bias=7)

    def in_sec(t):
        return 1e-7 < ((t - a0) % (2*np.pi)) < ad - 1e-7

    # Fuel rods + heat pipes — small circles with thin edges
    for p in FUEL_ALL:
        t = np.arctan2(p[1], p[0])
        if not in_sec(t):
            continue
        top_poly(circle_poly(p[0], p[1], FUEL_DO/2, n=16),
                 Z_TOP, C_FUEL, ec=EDGE, lw=0.18, bias=8)
    for p in HP_ALL:
        t = np.arctan2(p[1], p[0])
        if not in_sec(t):
            continue
        top_poly(circle_poly(p[0], p[1], HP_D/2, n=16),
                 Z_TOP, C_HP_STRIPE, ec=EDGE, lw=0.18, bias=8)

    # Control drums
    for dc in drum_centers():
        tdc = np.arctan2(dc[1], dc[0])
        dang = np.arcsin(min(1, DRUM_R / np.linalg.norm(dc)))
        pad = dang + np.deg2rad(2)
        if not (in_sec(tdc) and in_sec(tdc - pad) and in_sec(tdc + pad)):
            continue
        top_poly(circle_poly(dc[0], dc[1], DRUM_R, n=72),
                 Z_TOP, C_DRUM, ec=EDGE, lw=LW_MED, bias=9)
        # B4C crescent on OUTER side of drum — matches OpenMC:
        # drum_inner is offset by -2·cos(controldrum_angle) TOWARD the
        # core at angle=0, so the B4C annular region sits on the far
        # (outer) face of the drum.
        outward = np.arctan2(dc[1], dc[0])
        arc_a = np.linspace(outward - np.pi/3, outward + np.pi/3, 40)
        rdc_in = DRUM_R * 0.78
        cr = [[dc[0]+DRUM_R*np.cos(a), dc[1]+DRUM_R*np.sin(a)] for a in arc_a]
        cr += [[dc[0]+rdc_in*np.cos(a), dc[1]+rdc_in*np.sin(a)]
               for a in arc_a[::-1]]
        top_poly(cr, Z_TOP, C_B4C, ec=EDGE, lw=LW_MED, bias=10)


# -------- Cylinder wall --------------------------------------------
def draw_side(ca_deg, cb_deg, keep_wedge=False, n_strips=80):
    ca = np.deg2rad(ca_deg); cb = np.deg2rad(cb_deg)
    wd = (cb - ca) % (2*np.pi)
    if keep_wedge:
        vs = ca
        ve = ca + wd
        n_strips = max(14, int(round(n_strips * wd / (2*np.pi))) + 6)
    else:
        vs = cb
        ve = cb + (2*np.pi - wd)
    m = np.deg2rad(0.25)
    arc = np.linspace(vs + m, ve - m, n_strips + 1)
    for i in range(n_strips):
        a0, a1 = arc[i], arc[i+1]
        pts = [[R_REFL*np.cos(a0), R_REFL*np.sin(a0), Z_TOP],
               [R_REFL*np.cos(a1), R_REFL*np.sin(a1), Z_TOP],
               [R_REFL*np.cos(a1), R_REFL*np.sin(a1), Z_BOT],
               [R_REFL*np.cos(a0), R_REFL*np.sin(a0), Z_BOT]]
        v = np.array(pts)
        sx, sy, d = proj(v)
        # Flat fill, no stripes
        _push(d.mean(), np.column_stack([sx, sy]), C_SIDE, C_SIDE, 0.0)
    # Vertical silhouette edges at start/end of the arc
    for a in (arc[0], arc[-1]):
        pts = [[R_REFL*np.cos(a), R_REFL*np.sin(a), Z_TOP],
               [R_REFL*np.cos(a), R_REFL*np.sin(a), Z_BOT]]
        v = np.array(pts)
        sx, sy, d = proj(v)
        _push(d.mean() + 50,
              np.column_stack([sx, sy]), 'none', EDGE, LW_HEAVY)


# -------- Bottom face ----------------------------------------------
def draw_bottom(ca_deg, cb_deg, keep_wedge=False):
    ca = np.deg2rad(ca_deg); cb = np.deg2rad(cb_deg)
    wd = (cb - ca) % (2*np.pi)
    if keep_wedge:
        vs = ca
        ve = ca + wd
    else:
        vs = cb
        ve = cb + (2*np.pi - wd)
    m = np.deg2rad(0.3)
    n = max(24, int((ve - vs) * 60 / np.pi))
    arc = np.linspace(vs + m, ve - m, n)
    poly = [[0, 0, Z_BOT]] + [[R_REFL*np.cos(a), R_REFL*np.sin(a), Z_BOT] for a in arc]
    v = np.array(poly)
    sx, sy, _ = proj(v)
    _push(-1e5, np.column_stack([sx, sy]), np.clip(C_SIDE*0.78, 0, 1), EDGE, LW_MED)


# -------- Cut plane -------------------------------------------------
def cut_visible(theta, wc, is_wedge):
    d = (wc - theta + np.pi) % (2*np.pi) - np.pi
    na = theta + (np.pi/2 if d > 0 else -np.pi/2)
    if is_wedge:
        na += np.pi
    view = np.deg2rad(_VIEW['azim'])
    return np.cos(na - view) > 0


def draw_cut_plane(theta, wc, is_wedge=False):
    if not cut_visible(theta, wc, is_wedge):
        return
    c, s = np.cos(theta), np.sin(theta)

    def Q(r1, r2, z1, z2, color, bias=-5, lw=LW_THIN):
        pts = [[r1*c, r1*s, z1], [r2*c, r2*s, z1],
               [r2*c, r2*s, z2], [r1*c, r1*s, z2]]
        quad3d(pts, color, ec=EDGE, lw=lw, bias=bias)

    # Hex radial boundaries at this cut angle — use hex_dist so the
    # vertical cut-plane wall exactly meets the top-face hex polygon.
    r_outer = hex_dist(WALL_2, theta)
    r_inner = hex_dist(WALL_1, theta)

    # Reflector column (hex_outer → R_REFL)
    Q(r_outer, R_REFL, Z_BOT, Z_TOP, C_REFL, lw=LW_MED)
    # Central channel column (0 → hex_inner) from top of mauve band to crown
    Q(0, r_inner, Z_ABS_BOT, Z_TOP, C_CHANNEL)
    # Monolith cut face:
    #   body      (Z_FUEL_BOT → Z_FUEL_TOP)  solid salmon, no internals
    #   top band  (Z_FUEL_TOP → Z_TOP)       NO slab — only HP pin columns
    #                                         visible, background shows through
    Q(r_inner, r_outer, Z_FUEL_BOT, Z_FUEL_TOP, C_FUEL_SLAB, lw=LW_MED)
    Q(r_inner, r_outer, Z_ABS_BOT,  Z_ABS_TOP,  C_ABS_BAND)
    Q(0, r_outer, Z_PINK_BOT, Z_ABS_BOT, C_PINK_BAND)

    # 3D-shaded rod bundles in the top band.  Each rod is rendered as
    # N_SUB vertical slices with a cosine brightness profile across the
    # diameter so they read as cylindrical bundles (fuel + HP) rather
    # than flat rectangular bars.
    MARGIN  = 2.0
    hp_hw   = HP_D / 2.0
    fuel_hw = FUEL_DO / 2.0
    r_hp   = np.linspace(r_inner + MARGIN + hp_hw,
                         r_outer - MARGIN - hp_hw, 9)
    r_fuel = 0.5 * (r_hp[:-1] + r_hp[1:])
    N_SUB = 9

    def _shade(col, f):
        f = float(np.clip(f, 0.0, 1.6))
        out = np.clip(np.asarray(col, float) * f, 0.0, 1.0)
        return out

    def _rod(r_center, hw, base_color, bias):
        edges = np.linspace(-hw, +hw, N_SUB + 1)
        for i in range(N_SUB):
            u0, u1 = edges[i], edges[i+1]
            uc = 0.5 * (u0 + u1)
            bright = 0.50 + 0.65 * np.cos(0.5 * np.pi * uc / hw)
            col = _shade(base_color, bright)
            is_rim = (i == 0) or (i == N_SUB - 1)
            # Inner slices have no edge line (continuous gradient);
            # rim slices carry the silhouette stroke.
            if is_rim:
                pts = [[(r_center+u0)*c, (r_center+u0)*s, Z_FUEL_TOP],
                       [(r_center+u1)*c, (r_center+u1)*s, Z_FUEL_TOP],
                       [(r_center+u1)*c, (r_center+u1)*s, Z_TOP],
                       [(r_center+u0)*c, (r_center+u0)*s, Z_TOP]]
                quad3d(pts, col, ec=_shade(base_color, 0.40),
                       lw=0.25, bias=bias + i * 0.01)
            else:
                pts = [[(r_center+u0)*c, (r_center+u0)*s, Z_FUEL_TOP],
                       [(r_center+u1)*c, (r_center+u1)*s, Z_FUEL_TOP],
                       [(r_center+u1)*c, (r_center+u1)*s, Z_TOP],
                       [(r_center+u0)*c, (r_center+u0)*s, Z_TOP]]
                quad3d(pts, col, ec=col, lw=0.0,
                       bias=bias + i * 0.01)

    for r in r_hp:
        _rod(r, hp_hw, C_HP_STRIPE, bias=500)
    for r in r_fuel:
        _rod(r, fuel_hw, C_FUEL, bias=495)


# -------- Lattice generation (same as main file) -------------------
def gen_hp():
    p = []
    for i in range(4, 13):
        X = np.arange(WALL_1 + HP_D/2, 1000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_HP/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def gen_fuel():
    p = []
    for i in range(3, 11):
        X = np.arange(WALL_1 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5 + P_FUEL*np.sqrt(3)/2, -1000,
                      -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-3], Y[i-3] + j*P_HP])
    for i in range(4, 12):
        X = np.arange(WALL_1 + P_FUEL/2 + HP_D/2 + P_FUEL/2, 10000, P_FUEL*1.5)
        Y = np.arange(-P_HP*1.5, -1000, -P_FUEL*np.sqrt(3)/2)
        for j in range(i):
            p.append([X[i-4], Y[i-4] + j*P_HP])
    return np.array(p)


def replicate(pts):
    o = []
    for k in range(6):
        a = np.deg2rad(60*k)
        c, s = np.cos(a), np.sin(a)
        o.append(pts @ np.array([[c, -s], [s, c]]).T)
    return np.vstack(o)


def drum_centers():
    return [np.array([DRUM_CR*np.cos(np.deg2rad(60*k)),
                      DRUM_CR*np.sin(np.deg2rad(60*k))]) for k in range(6)]


HP_ALL = replicate(gen_hp())
FUEL_ALL = replicate(gen_fuel())


# -------- Panel ------------------------------------------------------
def render_panel(ax, ca_deg, cb_deg, *, azim, elev=25.0, is_wedge=False):
    global _Q
    _Q = []
    set_view(azim, elev)
    ca = np.deg2rad(ca_deg); cb = np.deg2rad(cb_deg)
    wd = (cb - ca) % (2*np.pi)
    wc = ca + wd/2
    draw_bottom(ca_deg, cb_deg, keep_wedge=is_wedge)
    draw_side(ca_deg, cb_deg, keep_wedge=is_wedge)
    draw_cut_plane(ca, wc, is_wedge=is_wedge)
    draw_cut_plane(cb, wc, is_wedge=is_wedge)
    draw_top_face(ca_deg, cb_deg, keep_wedge=is_wedge)
    _Q.sort(key=lambda f: f[0])
    for _, v, fc, ec, lw in _Q:
        is_none = isinstance(fc, str) and fc == 'none'
        if is_none:
            ax.add_patch(MplPoly(v, closed=False, fc='none', ec=ec, lw=lw))
        else:
            ax.add_patch(MplPoly(v, closed=True, fc=fc, ec=ec, lw=lw))
    ax.set_aspect('equal')
    ax.axis('off')
    ax.autoscale_view()
    return list(_Q)


# -------- Figure -----------------------------------------------------
fig, (axL, axR) = plt.subplots(1, 2, figsize=(20, 12),
                               gridspec_kw={'width_ratios': [1.6, 1.0]})

# LEFT: 4/6 cylinder (remove 120° centred at 210°, cut planes on clean
# 1/6 boundaries at 150° and 270°, each midway between adjacent drums).
LEFT_AZIM = 210; LEFT_CA = 150; LEFT_CB = 270
# RIGHT: axisymmetric 1/6 wedge 210°–270° containing drum at 240°.
# View azimuth mirrored (CB+60) so the dominant cut face swaps sides.
RIGHT_CA = 210;  RIGHT_CB = 270;  RIGHT_AZIM = (RIGHT_CA - 60) % 360
left_q  = render_panel(axL, LEFT_CA, LEFT_CB, azim=LEFT_AZIM, is_wedge=False)
right_q = render_panel(axR, RIGHT_CA, RIGHT_CB, azim=RIGHT_AZIM,
                       is_wedge=True)

plt.tight_layout(pad=1.0)

out = Path(__file__).parent
for fmt in ['svg', 'pdf', 'png', 'jpg']:
    p = out / f'core_cad_render.{fmt}'
    fig.savefig(p, format='jpeg' if fmt == 'jpg' else fmt,
                dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    print(f"Saved: {p}")


# -------- PPTX (draggable freeform shapes) ---------------------------
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor


def _rgb(c):
    if isinstance(c, str):
        if c == 'none':
            return None
        s = c.lstrip('#')
        if len(s) == 3:
            s = ''.join(ch*2 for ch in s)
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    r, g, b = [int(round(max(0, min(1, float(v))) * 255)) for v in c[:3]]
    return RGBColor(r, g, b)


def _bbox(q):
    all_v = np.concatenate([v for _, v, *_ in q], axis=0)
    return all_v[:, 0].min(), all_v[:, 1].min(), all_v[:, 0].max(), all_v[:, 1].max()


def _draw_freeforms(sl, q, bl, bt, bw, bh):
    x0, y0, x1, y1 = _bbox(q)
    pw, ph = x1 - x0, y1 - y0
    scale = min(bw / pw, bh / ph) * 0.97
    off_x = bl + (bw - pw * scale) / 2
    off_y = bt + (bh - ph * scale) / 2
    for _, v, fc, ec, lw in q:
        is_none_fc = isinstance(fc, str) and fc == 'none'
        is_none_ec = isinstance(ec, str) and ec == 'none'
        pts = [(Emu(int(off_x + (p[0]-x0)*scale)),
                Emu(int(off_y + (y1-p[1])*scale))) for p in v]
        if len(pts) < 2:
            continue
        ff = sl.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
        ff.add_line_segments(pts[1:], close=(not is_none_fc))
        shp = ff.convert_to_shape()
        if is_none_fc:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(fc)
        line = shp.line
        if lw and lw > 0:
            line.color.rgb = _rgb(ec) if not is_none_ec else RGBColor(0, 0, 0)
            line.width = Pt(max(0.25, float(lw)))
        else:
            line.fill.background()


prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
margin = Inches(0.25)
total_w = prs.slide_width - 2 * margin
total_h = prs.slide_height - 2 * margin
left_w  = int(total_w * (1.6 / 2.6))
right_w = int(total_w * (1.0 / 2.6))
gap     = Inches(0.15)
_draw_freeforms(slide, left_q,
                int(margin), int(margin),
                int(left_w - gap // 2), int(total_h))
_draw_freeforms(slide, right_q,
                int(margin) + int(left_w) + int(gap // 2),
                int(margin),
                int(right_w - gap // 2), int(total_h))
pptx_path = out / 'core_cad_render.pptx'
prs.save(str(pptx_path))
print(f"Saved: {pptx_path}")
plt.close(fig)
