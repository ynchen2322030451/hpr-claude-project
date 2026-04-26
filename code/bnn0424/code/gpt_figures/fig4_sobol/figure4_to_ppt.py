#!/usr/bin/env python3
"""Export Figure 4 to PPTX with native editable elements.

Slide 1: Full composed figure (image) + native title + panel labels
Slide 2-3: Individual panel A/B as images
Slide 4: Panel C summary — FULLY NATIVE (boxes + text, no image)
Slide 5: Editable text layer
"""

import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE.parent / "_shared"))
import ppt_helper as ppt

OUTDIR = HERE / "outputs"

BLUE_DARK  = "#1B3A5C"
BLUE_WASH  = "#E8EFF5"
BLUE_PALE  = "#C5DAE9"
ACCENT_ORA = "#C4713B"
ACCENT_TAN = "#FDF5E6"
GRAY_300   = "#B0B0B0"
GRAY_500   = "#7A7A7A"
GRAY_700   = "#4A4A4A"
GRAY_900   = "#1A1A1A"


def build_panel_c_native(slide):
    """Build Panel C (dominant-factor separation) entirely in native shapes."""
    # Section title
    ppt.add_text_box(slide,
        "C.  Dominant-factor separation summary",
        left=0.7, top=1.20, width=10.0, height=0.40,
        font_size=16, bold=True, color=GRAY_900)

    # Stress pathway box (blue scheme)
    ppt.add_rounded_rect(slide, left=1.0, top=2.0, width=5.0, height=3.2,
                         fill_color=BLUE_WASH, border_color=BLUE_DARK,
                         border_width=1.5)
    ppt.add_text_box(slide, "Stress pathway",
                     left=1.2, top=2.10, width=4.5, height=0.35,
                     font_size=14, bold=True, color=BLUE_DARK)
    stress_lines = [
        "Dominant factor: E intercept",
        "S\u2081 \u2248 0.54\u20130.58 (first-order)",
        "S_T \u2248 0.60\u20130.65 (total-order)",
        "Elastic modulus governs coupled",
        "  steady-state peak stress",
        "Secondary: \u03b1_base (thermal",
        "  expansion, via coupling)",
    ]
    ppt.add_multiline_text_box(
        slide, stress_lines,
        left=1.3, top=2.55, width=4.4, height=2.5,
        font_size=11, color=GRAY_700, bullet=True)

    # k_eff pathway box (orange scheme)
    ppt.add_rounded_rect(slide, left=6.8, top=2.0, width=5.0, height=3.2,
                         fill_color=ACCENT_TAN, border_color=ACCENT_ORA,
                         border_width=1.5)
    ppt.add_text_box(slide, "k_eff pathway",
                     left=7.0, top=2.10, width=4.5, height=0.35,
                     font_size=14, bold=True, color=ACCENT_ORA)
    keff_lines = [
        "Dominant factor: \u03b1_base",
        "S\u2081 \u2248 0.74\u20130.78 (first-order)",
        "S_T \u2248 0.78\u20130.84 (total-order)",
        "Thermal expansion base drives",
        "  temperature-dependent k_eff",
        "Secondary: SS316_k_ref",
        "  (conductivity reference)",
    ]
    ppt.add_multiline_text_box(
        slide, keff_lines,
        left=6.9, top=2.55, width=4.6, height=2.5,
        font_size=11, color=GRAY_700, bullet=True)

    # Dashed separator line
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    ppt.add_line(slide, 6.3, 2.2, 6.3, 5.0,
                 color=GRAY_300, width=1.0,
                 dash=MSO_LINE_DASH_STYLE.DASH)

    # Bottom note
    ppt.add_text_box(
        slide,
        "Different material parameters dominate each output pathway, "
        "enabling targeted uncertainty reduction strategies.",
        left=1.0, top=5.50, width=10.5, height=0.50,
        font_size=10, italic=True, color=GRAY_500)


def main():
    prs = ppt.create_presentation()

    # ── Slide 1: Full composed figure ─────────────────────────────────────
    slide1 = ppt.add_blank_slide(prs)
    ppt.add_slide_title(slide1,
        "Figure 4 | Sobol variance decomposition",
        font_size=18, color=GRAY_900)
    ppt.add_centered_image(slide1, OUTDIR / "figure4_full.png",
                           max_w=12.0, max_h=5.8, y_offset=0.55)
    ppt.add_panel_label(slide1, "A", left=0.30, top=0.90)
    ppt.add_panel_label(slide1, "B", left=0.30, top=3.80)
    ppt.add_panel_label(slide1, "C", left=8.00, top=0.90)

    # ── Slide 2: Panel A — Stress Sobol ───────────────────────────────────
    img_a = OUTDIR / "panel_A_stress.png"
    if img_a.exists():
        slide2 = ppt.add_blank_slide(prs)
        ppt.add_slide_title(slide2,
            "Figure 4A | Coupled steady-state max stress — Sobol indices",
            font_size=18, color=GRAY_900)
        ppt.add_centered_image(slide2, img_a, max_w=10.0, max_h=5.8,
                               y_offset=0.55)
        ppt.add_panel_label(slide2, "A", left=0.50, top=0.90)

    # ── Slide 3: Panel B — k_eff Sobol ────────────────────────────────────
    img_b = OUTDIR / "panel_B_keff.png"
    if img_b.exists():
        slide3 = ppt.add_blank_slide(prs)
        ppt.add_slide_title(slide3,
            "Figure 4B | k_eff (coupled) — Sobol indices",
            font_size=18, color=GRAY_900)
        ppt.add_centered_image(slide3, img_b, max_w=10.0, max_h=5.8,
                               y_offset=0.55)
        ppt.add_panel_label(slide3, "B", left=0.50, top=0.90)

    # ── Slide 4: Panel C — FULLY NATIVE summary ──────────────────────────
    slide4 = ppt.add_blank_slide(prs)
    ppt.add_slide_title(slide4,
        "Figure 4C | Dominant-factor separation (native editable)",
        font_size=18, color=GRAY_900)
    build_panel_c_native(slide4)

    # ── Slide 5: Editable text layer ──────────────────────────────────────
    slide_txt = ppt.add_blank_slide(prs)
    ppt.add_slide_title(slide_txt,
        "Figure 4 — Editable text elements",
        font_size=16, color=GRAY_500)

    y = 1.20
    for letter, title in [
        ("A", "Coupled steady-state max stress — Sobol indices (S\u2081 / S_T)"),
        ("B", "k_eff (coupled) — Sobol indices (S\u2081 / S_T)"),
        ("C", "Dominant-factor separation summary"),
    ]:
        ppt.add_text_box(slide_txt, f"{letter}.  {title}",
                         left=0.7, top=y, width=10.0, height=0.35,
                         font_size=13, bold=True, color=BLUE_DARK)
        y += 0.45

    y += 0.30
    ppt.add_text_box(slide_txt, "Key findings (editable):",
                     left=0.7, top=y, width=8.0, height=0.30,
                     font_size=11, bold=True, color=GRAY_900)
    y += 0.35
    findings = [
        "E intercept dominates stress (S\u2081 \u2248 0.54\u20130.58)",
        "\u03b1_base dominates k_eff (S\u2081 \u2248 0.74\u20130.78)",
        "Distinct sensitivity pathways enable targeted uncertainty reduction",
        "95% bootstrap CI from 4 independent BNN chains",
    ]
    ppt.add_multiline_text_box(
        slide_txt, findings,
        left=0.9, top=y, width=11.0, height=1.5,
        font_size=10, color=GRAY_900, bullet=True)

    y += 1.8
    ppt.add_text_box(
        slide_txt,
        "Caption: Sobol variance decomposition reveals distinct dominant "
        "factors for stress (elastic modulus) and k_eff (thermal expansion), "
        "suggesting that targeted experimental characterization of these "
        "material properties would most effectively reduce response uncertainty.",
        left=0.7, top=y, width=11.0, height=1.0,
        font_size=10, italic=True, color=GRAY_500)

    notes = (
        "Figure 4: Sobol variance decomposition.\n"
        "Panel A: stress, Panel B: k_eff, Panel C: summary.\n"
        "Panel C on slide 4 is fully native — every box and text is editable."
    )
    prs.slides[0].notes_slide.notes_text_frame.text = notes

    ppt.save_pptx(prs, OUTDIR / "figure4_sobol.pptx")


if __name__ == "__main__":
    main()
