"""
Hybrid SVG→PPTX converter: PNG background + native editable text overlays.

For each figure:
  1. Embed the high-quality PNG as a full-slide background image
  2. Parse the SVG to extract ALL <text> elements (axis labels, tick values,
     titles, panel labels, legend text, annotations)
  3. Place each text element as a native PPT TextBox at the correct position

Result: visually identical to the original figure, but every text element
can be double-clicked and edited in PowerPoint.

Usage:
    from svg_to_hybrid_pptx import convert_figure, batch_convert

    # Single figure
    convert_figure("figure.svg", "figure.png", "output.pptx", title="Figure 1")

    # Batch
    batch_convert([
        ("fig1.svg", "fig1.png", "fig1.pptx", "Figure 1"),
        ...
    ])
"""
from __future__ import annotations

import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ── Constants ────────────────────────────────────────────────────────────
SVG_NS = "http://www.w3.org/2000/svg"
PT_TO_EMU = 12700
IN_TO_EMU = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ── Text element dataclass ──────────────────────────────────────────────

@dataclass
class TextElem:
    text: str
    x_pt: float      # position in SVG viewBox coords (pt)
    y_pt: float
    font_size: float  # pt
    font_family: str
    bold: bool
    italic: bool
    color: str        # hex '#rrggbb'
    rotation: float   # degrees
    anchor: str       # 'start', 'middle', 'end'


# ── SVG text extraction ─────────────────────────────────────────────────

def _local_tag(elem: ET.Element) -> str:
    tag = elem.tag
    return tag.split("}", 1)[1] if "}" in tag else tag


def _parse_transform(tf_str: str) -> Tuple[float, float, float]:
    """Parse SVG transform string, return (tx, ty, rotation_deg)."""
    tx, ty, rot = 0.0, 0.0, 0.0
    if not tf_str:
        return tx, ty, rot

    # translate(x, y)
    m = re.search(r"translate\(\s*([+-]?[\d.]+)[,\s]+([+-]?[\d.]+)\s*\)", tf_str)
    if m:
        tx, ty = float(m.group(1)), float(m.group(2))

    # rotate(deg)
    m = re.search(r"rotate\(\s*([+-]?[\d.]+)", tf_str)
    if m:
        rot = float(m.group(1))

    # scale — not common for text, but handle simple case
    m = re.search(r"scale\(\s*([+-]?[\d.]+)", tf_str)
    if m:
        # scale affects font size, handled separately
        pass

    return tx, ty, rot


def _get_inherited_style(elem: ET.Element, attr: str, default: str = "") -> str:
    """Walk up to find an attribute value (simple style inheritance)."""
    val = elem.get(attr, "")
    if val:
        return val
    # Check inline style
    style = elem.get("style", "")
    if style:
        for part in style.split(";"):
            if attr in part and ":" in part:
                k, v = part.split(":", 1)
                if k.strip() == attr:
                    return v.strip()
    return default


def extract_texts_from_svg(svg_path: Path) -> Tuple[float, float, List[TextElem]]:
    """Parse SVG and extract all text elements.

    Returns (svg_width_pt, svg_height_pt, list_of_TextElem).
    """
    tree = ET.parse(str(svg_path))
    root = tree.getroot()

    # Get SVG dimensions
    w_str = root.get("width", "")
    h_str = root.get("height", "")
    if w_str and h_str:
        svg_w = float(re.sub(r"[a-zA-Z]+", "", w_str))
        svg_h = float(re.sub(r"[a-zA-Z]+", "", h_str))
    else:
        vb = root.get("viewBox", "0 0 800 600")
        parts = vb.split()
        svg_w, svg_h = float(parts[2]), float(parts[3])

    texts: List[TextElem] = []

    # Walk all <text> elements
    for text_elem in root.iter(f"{{{SVG_NS}}}text"):
        t = _parse_text_element(text_elem)
        if t:
            texts.append(t)

    # Also walk <g> groups for transform-wrapped text
    # (already handled by iter finding nested <text>)

    return svg_w, svg_h, texts


def _parse_text_element(elem: ET.Element) -> Optional[TextElem]:
    """Parse a single <text> element, collecting tspan children."""
    # Collect text content
    parts = []
    if elem.text and elem.text.strip():
        parts.append(elem.text.strip())

    for child in elem:
        child_tag = _local_tag(child)
        if child_tag == "tspan":
            if child.text and child.text.strip():
                parts.append(child.text.strip())
        if child.tail and child.tail.strip():
            parts.append(child.tail.strip())

    text = " ".join(parts)
    if not text:
        return None

    # Position — x/y may be space-separated lists (one per char); take first
    def _first_float(s: str, default: float = 0.0) -> float:
        s = s.strip()
        if not s:
            return default
        return float(s.split()[0])

    x = _first_float(elem.get("x", "0"))
    y = _first_float(elem.get("y", "0"))

    # Check for tspan with its own x,y (use first tspan position if present)
    for child in elem:
        if _local_tag(child) == "tspan":
            tx_s = child.get("x", "")
            ty_s = child.get("y", "")
            if tx_s:
                x = _first_float(tx_s)
            if ty_s:
                y = _first_float(ty_s)
            break  # use first tspan

    # Transform on <text> element
    tf_str = elem.get("transform", "")
    tx, ty, rot = _parse_transform(tf_str)
    x += tx
    y += ty

    # Also check parent <g> transforms by walking up
    # (ElementTree doesn't provide parent access, so we handle
    #  the common case of transform on the <text> element itself)

    # Font size
    font_size = 10.0
    fs_str = elem.get("font-size", "")
    if not fs_str:
        style = elem.get("style", "")
        m = re.search(r"font-size:\s*([\d.]+)", style)
        if m:
            fs_str = m.group(1)
    if fs_str:
        try:
            font_size = float(re.sub(r"[a-zA-Z]+", "", fs_str))
        except ValueError:
            pass

    # Font family
    family = elem.get("font-family", "")
    if not family:
        style = elem.get("style", "")
        m = re.search(r"font-family:\s*([^;]+)", style)
        if m:
            family = m.group(1)
    family = family.strip("'\" ")
    if not family or "dejavu" in family.lower() or "sans" in family.lower():
        family = "Arial"
    if "serif" in family.lower() and "sans" not in family.lower():
        family = "Times New Roman"

    # Bold
    weight = elem.get("font-weight", "")
    if not weight:
        style = elem.get("style", "")
        m = re.search(r"font-weight:\s*(\w+)", style)
        if m:
            weight = m.group(1)
    bold = weight.lower() in ("bold", "700", "800", "900")

    # Italic
    fstyle = elem.get("font-style", "")
    if not fstyle:
        style = elem.get("style", "")
        m = re.search(r"font-style:\s*(\w+)", style)
        if m:
            fstyle = m.group(1)
    italic = fstyle.lower() in ("italic", "oblique")

    # Color (fill attribute)
    color = "#000000"
    fill = elem.get("fill", "")
    if fill.startswith("#"):
        color = fill
    else:
        style = elem.get("style", "")
        m = re.search(r"(?:^|;)\s*fill:\s*(#[0-9a-fA-F]{3,6})", style)
        if m:
            color = m.group(1)

    # Text anchor
    anchor = elem.get("text-anchor", "start")
    if not anchor or anchor == "inherit":
        style = elem.get("style", "")
        m = re.search(r"text-anchor:\s*(\w+)", style)
        anchor = m.group(1) if m else "start"

    return TextElem(
        text=text, x_pt=x, y_pt=y,
        font_size=font_size, font_family=family,
        bold=bold, italic=italic,
        color=color, rotation=rot, anchor=anchor,
    )


# ── PPTX generation ─────────────────────────────────────────────────────

def _hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def convert_figure(svg_path: str, png_path: str, output_path: str,
                   title: str = "") -> str:
    """Convert one figure to hybrid PPTX.

    Args:
        svg_path: Path to SVG (for text extraction)
        png_path: Path to PNG (for background image)
        output_path: Output PPTX path
        title: Optional slide title (added as extra TextBox)

    Returns output_path.
    """
    svg_path = Path(svg_path)
    png_path = Path(png_path)

    if not png_path.exists():
        raise FileNotFoundError(f"PNG not found: {png_path}")

    # Extract text elements from SVG
    if svg_path.exists():
        svg_w, svg_h, texts = extract_texts_from_svg(svg_path)
    else:
        svg_w, svg_h, texts = 100, 100, []

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Add PNG as background image (centered, fill slide)
    img = Image.open(png_path)
    iw, ih = img.size
    img.close()
    aspect = iw / ih

    margin = 0.3  # inches
    max_w = SLIDE_W_IN - 2 * margin
    max_h = SLIDE_H_IN - 2 * margin

    if aspect > max_w / max_h:
        pic_w = max_w
        pic_h = pic_w / aspect
    else:
        pic_h = max_h
        pic_w = pic_h * aspect

    pic_left = (SLIDE_W_IN - pic_w) / 2
    pic_top = (SLIDE_H_IN - pic_h) / 2

    pic = slide.shapes.add_picture(
        str(png_path),
        Inches(pic_left), Inches(pic_top),
        Inches(pic_w), Inches(pic_h))

    # Coordinate mapping: SVG pt → slide inches
    # The PNG occupies [pic_left, pic_top] to [pic_left+pic_w, pic_top+pic_h]
    # SVG viewBox is [0, 0, svg_w, svg_h]
    scale_x = pic_w / svg_w  # inches per SVG-pt
    scale_y = pic_h / svg_h
    offset_x = pic_left  # inches
    offset_y = pic_top

    # Place text overlays
    n_placed = 0
    for t in texts:
        # Map SVG position to slide position (inches)
        x_in = t.x_pt * scale_x + offset_x
        y_in = t.y_pt * scale_y + offset_y

        # Scale font size
        fs_pt = t.font_size * scale_y  # scale by vertical ratio
        fs_pt = max(fs_pt, 4.0)  # minimum readable

        # Estimate text box dimensions
        char_w = fs_pt * 0.55 / 72 * 96 / 72  # rough char width in inches
        # Simpler: use font_size in pt → width in inches
        char_w_in = fs_pt * 0.012  # empirical: ~0.012 inches per pt per char
        box_w = max(len(t.text) * char_w_in + 0.1, 0.3)
        box_h = fs_pt / 72 * 1.6  # font height + padding

        # Adjust position based on text-anchor
        if t.anchor == "middle":
            x_in -= box_w / 2
        elif t.anchor == "end":
            x_in -= box_w

        # Baseline to top adjustment
        y_in -= fs_pt / 72

        # Handle rotation
        is_rotated = abs(t.rotation) > 1.0

        # Create TextBox
        txbox = slide.shapes.add_textbox(
            Inches(x_in), Inches(y_in),
            Inches(box_w), Inches(box_h))

        if is_rotated:
            rot = t.rotation
            if rot < 0:
                rot += 360.0
            txbox.rotation = rot

        tf = txbox.text_frame
        tf.word_wrap = False
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p = tf.paragraphs[0]
        if t.anchor == "middle":
            p.alignment = PP_ALIGN.CENTER
        elif t.anchor == "end":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = t.text
        font = run.font
        font.size = Pt(fs_pt)
        font.name = t.font_family
        font.bold = t.bold
        font.italic = t.italic
        if t.color and t.color != "none":
            try:
                font.color.rgb = _hex_to_rgb(t.color)
            except Exception:
                pass

        n_placed += 1

    prs.save(output_path)
    return output_path


def batch_convert(items: list, verbose: bool = True) -> List[str]:
    """Convert multiple figures.

    items: list of (svg_path, png_path, output_path, title)
    Returns list of output paths.
    """
    outputs = []
    for svg_path, png_path, output_path, title in items:
        svg_p = Path(svg_path)
        png_p = Path(png_path)
        out_p = Path(output_path)

        if not png_p.exists():
            if verbose:
                print(f"  [skip] {png_p.name} (PNG not found)")
            continue

        try:
            n_texts = 0
            if svg_p.exists():
                _, _, texts = extract_texts_from_svg(svg_p)
                n_texts = len(texts)

            convert_figure(str(svg_p), str(png_p), str(out_p), title=title)
            outputs.append(str(out_p))

            if verbose:
                print(f"  [ok]   {out_p.parent.name}/{out_p.name}  "
                      f"({n_texts} editable text elements)")
        except Exception as e:
            if verbose:
                print(f"  [FAIL] {out_p.name}: {e}")

    return outputs
