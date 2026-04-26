"""
PPT export helper — native editable elements.

All shapes, text boxes, and connectors are native PowerPoint objects:
  - Every text box is editable (font, size, color, position)
  - Every shape is movable and recolorable
  - Connectors/arrows can be repositioned

Provides building blocks used by each figure's _to_ppt.py script.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Slide dimensions ──────────────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5


# ── Color utilities ───────────────────────────────────────────────────────
def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' to pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Presentation factory ─────────────────────────────────────────────────
def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def save_pptx(prs: Presentation, path):
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    print(f"  Saved PPTX: {path.name} ({len(prs.slides)} slides)")


# ── Image insertion ───────────────────────────────────────────────────────
def add_centered_image(slide, img_path, max_w=12.0, max_h=6.0, y_offset=0.4):
    """Embed a PNG/JPEG centered on the slide. Returns the picture shape."""
    from PIL import Image
    img_path = Path(img_path)
    if not img_path.exists():
        print(f"  Warning: {img_path.name} not found, skipping")
        return None
    img = Image.open(img_path)
    w_px, h_px = img.size
    aspect = w_px / h_px
    if aspect > max_w / max_h:
        w = max_w
        h = w / aspect
    else:
        h = max_h
        w = h * aspect
    left = (SLIDE_W - w) / 2
    top = (SLIDE_H - h) / 2 + y_offset
    return slide.shapes.add_picture(
        str(img_path), Inches(left), Inches(top), Inches(w), Inches(h))


# ── Native text boxes ────────────────────────────────────────────────────
def add_slide_title(slide, text, font_size=20, bold=True,
                    color="#1A1A1A", left=0.5, top=0.15, width=12.0, height=0.55):
    """Add an editable title text box at top of slide."""
    return add_text_box(slide, text, left, top, width, height,
                        font_size=font_size, bold=bold, color=color,
                        alignment=PP_ALIGN.LEFT)


def add_text_box(slide, text, left, top, width, height,
                 font_size=12, bold=False, italic=False,
                 color="#1A1A1A", alignment=PP_ALIGN.LEFT,
                 font_name="Arial", vertical_anchor=MSO_ANCHOR.TOP):
    """Add a native editable text box at specified position (inches)."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    if vertical_anchor:
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)
    run.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text_box(slide, lines, left, top, width, height,
                           font_size=10, color="#4A4A4A",
                           font_name="Arial", line_spacing=1.2,
                           bullet=False):
    """Add a text box with multiple paragraphs (one per line)."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        display = f"\u2022  {line}" if bullet else line
        p.text = display
        p.space_after = Pt(2)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = hex_to_rgb(color)
            run.font.name = font_name
    return txBox


def add_panel_label(slide, letter, left, top, font_size=18):
    """Add a bold panel label (A, B, C ...) as a native text box."""
    return add_text_box(slide, letter, left, top, 0.4, 0.35,
                        font_size=font_size, bold=True, color="#1A1A1A")


# ── Native shapes ────────────────────────────────────────────────────────
def add_rounded_rect(slide, left, top, width, height,
                     fill_color="#E8EFF5", border_color="#B0B0B0",
                     border_width=1.0, corner_radius=None):
    """Add a native rounded rectangle. Returns the shape object."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.color.rgb = hex_to_rgb(border_color)
    shape.line.width = Pt(border_width)
    # Adjust corner rounding if specified (0.0 = square, 1.0 = pill)
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_rect(slide, left, top, width, height,
             fill_color="#E8EFF5", border_color=None, border_width=0):
    """Add a native rectangle (no rounding)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    if border_color:
        shape.line.color.rgb = hex_to_rgb(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color="#7A7A7A", width=1.5, dash=None):
    """Add a native line shape."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE as MCT
    connector = slide.shapes.add_connector(
        MCT.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        connector.line.dash_style = dash
    return connector


def add_arrow(slide, x1, y1, x2, y2, color="#7A7A7A", width=1.5):
    """Add a native arrow connector (line with arrowhead at end)."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE as MCT
    connector = slide.shapes.add_connector(
        MCT.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = hex_to_rgb(color)
    connector.line.width = Pt(width)
    # Add arrowhead via XML — access <a:ln> inside <p:spPr>
    spPr = connector._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = connector.line._get_or_add_ln()
    from lxml import etree
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("len", "med")
    return connector


# ── Box with text frame (for workflow-type diagrams) ──────────────────────
def add_box_with_text(slide, left, top, width, height,
                      title, body_lines,
                      fill_color="#E8EFF5", border_color="#B0B0B0",
                      title_color="#3B6B9A", body_color="#4A4A4A",
                      title_size=12, body_size=9.5, border_width=1.0,
                      corner_radius=None):
    """Add a rounded rectangle with title + bullet body text inside.

    Returns the shape (text is part of the shape's text_frame).
    """
    shape = add_rounded_rect(
        slide, left, top, width, height,
        fill_color=fill_color, border_color=border_color,
        border_width=border_width, corner_radius=corner_radius)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.06)

    # Title paragraph
    p = tf.paragraphs[0]
    p.text = title
    p.space_after = Pt(6)
    run = p.runs[0]
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = hex_to_rgb(title_color)
    run.font.name = "Arial"

    # Body bullet lines
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = f"\u2022  {line}"
        p.space_after = Pt(2)
        for run in p.runs:
            run.font.size = Pt(body_size)
            run.font.color.rgb = hex_to_rgb(body_color)
            run.font.name = "Arial"

    return shape


# ── High-level: single-image slide with native title ─────────────────────
def image_slide(prs, img_path, title, max_w=12.0, max_h=6.0):
    """Create a slide with centered image + native editable title.

    Returns (slide, picture_shape).
    """
    slide = add_blank_slide(prs)
    add_slide_title(slide, title)
    pic = add_centered_image(slide, img_path, max_w=max_w, max_h=max_h)
    return slide, pic


# ── Legacy compatibility ─────────────────────────────────────────────────
def png_to_pptx(png_path, pptx_path, title="", notes="",
                max_w_inches=12.0, max_h_inches=6.5):
    """Create a single-slide PPTX with centered image + native title."""
    prs = create_presentation()
    slide, _ = image_slide(prs, png_path, title,
                           max_w=max_w_inches, max_h=max_h_inches)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    save_pptx(prs, pptx_path)


def multi_panel_pptx(image_paths, pptx_path, title="", notes=""):
    """Create PPTX with multiple images, one per slide, native titles."""
    prs = create_presentation()
    for img_path in image_paths:
        img_path = Path(img_path)
        if not img_path.exists():
            print(f"  Warning: {img_path} not found, skipping")
            continue
        slide_title = f"{title} — {img_path.stem}" if title else img_path.stem
        slide, _ = image_slide(prs, img_path, slide_title)

    if notes and len(prs.slides) > 0:
        prs.slides[0].notes_slide.notes_text_frame.text = notes
    save_pptx(prs, pptx_path)
